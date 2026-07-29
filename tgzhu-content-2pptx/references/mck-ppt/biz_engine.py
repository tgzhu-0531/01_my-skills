# -*- coding: utf-8 -*-
"""
Biz Engine — 混合方案B
模板克隆（保留品牌区） + 内容布局 + 企业配色
基于 Mck-ppt-design-skill 的布局算法适配
"""
import os, shutil, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══════════════════════════════════════════
# 企业品牌色
# ═══════════════════════════════════════════
BIZ_PRIMARY = RGBColor(0x00, 0xA7, 0xCB)
BIZ_DARK    = RGBColor(0x38, 0x4F, 0x63)
# 企业风：内容卡底用深蓝（与模板的深色协调），文字用白色
BIZ_CARD_BG = RGBColor(0x1A, 0x2D, 0x42)   # 深蓝卡底
BIZ_CARD_BORDER = RGBColor(0x00, 0xA7, 0xCB)  # 品牌青边框
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT   = RGBColor(0xE5, 0xE5, 0xE5)
GOLD         = RGBColor(0xC9, 0xA2, 0x27)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MED_GRAY     = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY      = RGBColor(0xF2, 0xF2, 0xF2)
NAVY         = RGBColor(0x05, 0x1C, 0x2C)
BIZ_GOLD    = RGBColor(0xC9, 0xA2, 0x27)
BIZ_ACCENT  = RGBColor(0x00, 0xAF, 0xD2)
BIZ_ORANGE  = RGBColor(0xF7, 0x96, 0x46)

# ═══════════════════════════════════════════
# 布局常量（来自 Mck Engine）
# ═══════════════════════════════════════════
SW = Inches(13.333)
SH = Inches(7.5)
LM = Inches(0.8)
RM = Inches(0.8)
CW = Inches(11.733)

TITLE_TOP       = Inches(0.15)
TITLE_H         = Inches(0.9)
TITLE_LINE_Y    = Inches(1.05)
CONTENT_TOP     = Inches(1.3)
SOURCE_Y        = Inches(7.05)
PAGE_NUM_X      = Inches(12.2)
BOTTOM_BAR_Y    = Inches(6.2)
BOTTOM_BAR_H    = Inches(0.65)

COVER_TITLE_SIZE   = Pt(44)
ACTION_TITLE_SIZE  = Pt(22)
SUB_HEADER_SIZE    = Pt(18)
BODY_SIZE          = Pt(14)
SMALL_SIZE         = Pt(12)
FOOTNOTE_SIZE      = Pt(9)

# ═══════════════════════════════════════════
# 核心绘图函数（来自 Mck core.py）
# ═══════════════════════════════════════════

def _clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

def set_ea_font(run, typeface='微软雅黑'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)

def add_text(slide, left, top, width, height, text,
             font_size=BODY_SIZE, font_name='Arial',
             font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='微软雅黑',
             anchor=MSO_ANCHOR.TOP, line_spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = font_size; p.font.name = font_name
        p.font.color.rgb = font_color; p.font.bold = bold; p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0); p.space_after = Pt(0)
        p.line_spacing = 0.93 if font_size.pt >= 18 else Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); _clean_shape(shape)
    return shape

def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)

def add_action_title(slide, text):
    # 标题在顶部左侧：限制宽度避开右上角品牌区（Logo+slogan）
    # 模板品牌区从 x≈9in 开始 → 标题只到 x≈9.2in
    add_text(slide, Inches(0.46), Inches(0.25), Inches(8.6), Inches(0.65), text,
             font_size=Pt(22), font_color=WHITE, bold=True,
             font_name='微软雅黑', anchor=MSO_ANCHOR.BOTTOM)

def add_source(slide, text, y=SOURCE_Y):
    # 来源行在深色背景上 → 用浅灰
    add_text(slide, LM, y, Inches(11), Inches(0.3), text,
             font_size=FOOTNOTE_SIZE, font_color=RGBColor(0xCC, 0xCC, 0xCC))

def add_page_number(slide, num, total):
    # 页码在深色背景上 → 用浅灰
    add_text(slide, PAGE_NUM_X, Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=FOOTNOTE_SIZE, font_color=RGBColor(0xCC, 0xCC, 0xCC),
             alignment=PP_ALIGN.RIGHT)

def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

# ═══════════════════════════════════════════
# BizEngine
# ═══════════════════════════════════════════

class BizEngine:
    """基于模板克隆的混合布局引擎。"""
    
    def __init__(self, template_path, output_path, total_slides=30):
        self.template_path = template_path
        self.output_path = output_path
        self.total = total_slides
        self._page = 0
        self._cover_used = False  # 封面是否已生成（用模板自带第一页）
        shutil.copy2(template_path, output_path)
        self.prs = Presentation(output_path)
        # 不删除默认页，用它作为封面
    
    def save(self):
        self.prs.save(self.output_path)
        full_cleanup(self.output_path)
        print(f"✅ 已生成: {self.output_path} ({self._page}页)")
    
    def _content_slide(self):
        self._page += 1
        return self.prs.slides.add_slide(self.prs.slide_masters[0].slide_layouts[6])
    
    # ─── 封面（用模板自带第一页）────────────
    def cover(self, title, subtitle=''):
        self._page += 1
        s = self.prs.slides[0]  # 使用模板自带的第一页
        for sh in s.shapes:
            if not hasattr(sh, 'text_frame'): continue
            txt = sh.text.strip()
            if "主标题" in txt:
                tf = sh.text_frame; tf.clear()
                for a in ['margin_left','margin_right','margin_top','margin_bottom']:
                    setattr(tf, a, Emu(0))
                p = tf.paragraphs[0]; p.font.bold = True
                r = p.add_run(); r.text = title
                r.font.size = Pt(29.65); r.font.bold = True
                r.font.color.rgb = WHITE; r.font.name = '微软雅黑'
                set_ea_font(r, '微软雅黑')
            elif "副标题" in txt:
                sh.text_frame.clear()
        if subtitle:
            tb = s.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.5))
            tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            for a in ['margin_left','margin_right','margin_top','margin_bottom']:
                setattr(tf, a, Emu(0))
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = subtitle
            r.font.size = Pt(36); r.font.bold = True
            r.font.color.rgb = BIZ_ACCENT; r.font.name = '微软雅黑'
            set_ea_font(r, '微软雅黑')
        print(f"  ✓ 封面: {title}")
        return s
    
    # ─── 内容页（标准） ────────────────────
    def content(self, title, subtitle='', lines=None, source=''):
        s = self._content_slide()
        add_action_title(s, title)
        # 副标题在标题下方
        sub_top = Inches(0.95)
        if subtitle:
            add_text(s, Inches(0.46), sub_top, Inches(8.6), Inches(0.4), subtitle,
                     font_size=Pt(14), font_color=BIZ_PRIMARY, bold=True,
                     font_name='微软雅黑')
        # 内容区从模板品牌线下方开始
        y = Inches(1.55)
        if lines:
            # 深蓝卡底 + 左侧品牌青细线作装饰（去掉顶部分隔线）
            add_rect(s, Inches(0.46), y, Inches(12.4), Inches(4.8), BIZ_CARD_BG)
            add_rect(s, Inches(0.46), y, Pt(3), Inches(4.8), BIZ_CARD_BORDER)
            normal = [l for l in lines if not l.startswith('→')]
            highlight = [l for l in lines if l.startswith('→')]
            if normal:
                add_text(s, Inches(0.76), y + Inches(0.2), Inches(12.0), Inches(3.2),
                         normal, font_size=BODY_SIZE, font_color=LIGHT_TEXT,
                         font_name='微软雅黑', line_spacing=Pt(8))
            if highlight:
                hy = y + Inches(3.6)
                for i, hl in enumerate(highlight):
                    add_text(s, Inches(0.76), hy + Inches(0.5) * i, Inches(12.0), Inches(0.5),
                             hl, font_size=Pt(16), font_color=GOLD, bold=True,
                             font_name='微软雅黑')
        # 注意：不再添加底部来源和页码（按用户要求去掉）
        print(f"  ✓ {title}")
        return s
    
    # ─── 双栏对比页 ────────────────────────
    def two_column(self, title, col1_title, col1_items, col2_title, col2_items, source=''):
        s = self._content_slide()
        add_action_title(s, title)
        cw = Inches(6.0)
        cy = Inches(1.55)
        cx1 = Inches(0.46); cx2 = Inches(0.46) + cw + Inches(0.4)
        for cx, ct, ci in [(cx1, col1_title, col1_items), (cx2, col2_title, col2_items)]:
            # 卡底：左边缘一道品牌青细线（左侧标识）作装饰
            add_rect(s, cx, cy, cw, Inches(4.8), BIZ_CARD_BG)
            add_rect(s, cx, cy, Pt(3), Inches(4.8), BIZ_CARD_BORDER)
            # 列标题：品牌青大字号
            add_text(s, cx + Inches(0.3), cy + Inches(0.2), cw - Inches(0.4), Inches(0.5),
                     ct, font_size=Pt(18), font_color=BIZ_PRIMARY, bold=True,
                     font_name='微软雅黑')
            # 列表项：浅灰文字，每项前加圆点
            for i, item in enumerate(ci if isinstance(ci, list) else [ci]):
                add_text(s, cx + Inches(0.3), cy + Inches(0.85) + Inches(0.55) * i, cw - Inches(0.4), Inches(0.5),
                         f"·  {item}", font_size=BODY_SIZE, font_color=LIGHT_TEXT,
                         font_name='微软雅黑')
        # 注意：不再添加底部来源和页码
        print(f"  ✓ {title}")
        return s
    
    def _footer(self, s, source=None):
        if source: add_source(s, source)
        add_page_number(s, self._page, self.total)
