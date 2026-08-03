# -*- coding: utf-8 -*-
"""
business.py — 商务风引擎（McKinsey 咨询风格，白底深字）
不依赖任何模板：从空白 16:9 画布代码绘制，全部版式绝对坐标手绘。
视觉特征（白底深字，与原「克隆企业模板」彻底解耦）：

  - 白底（blank_deck 默认），行动标题（action title）置顶左、结论式、
    22pt 深蓝 BIZ_DARK 粗体（McKinsey 标志）
  - 内容卡：浅灰卡底 #F2F2F2 + 左侧品牌青竖线 #00A7CB（非深底白字）
  - 高亮行（hl / "→"）：金色 #C9A227
  - 全局中文字体：微软雅黑

暴露统一版式原语：cover / section_header / timeline / two_column / card_grid /
summary / bottom_gold （与 fengmang.py / enterprise.py / wenzhi.py 同签名，
调用方可用同一套 runner 驱动所有风格）。
"""
import os, zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from common import (BIZ_TEAL, BIZ_TEAL2, BIZ_DARK, WHITE, GOLD, LIGHT_GRAY,
                    SP_UNIT, SP_GROUP, SP_CARD_BODY, SP_RUN_DEF,
                    tb_box, run_text, rect, add_page_number, shape_bottom_in,
                    set_chinese_font, blank_deck, set_bg, FONT,
                    cover_qr, bottom_gold as _gold, cover_line, fit_cover_title,
                    layout_timeline, layout_card_grid, layout_two_column, layout_summary,
                    layout_compare, layout_loop_page, layout_profile_warning,
                    layout_transition_rows)

# 商务风专属色（白底深字：三级色彩体系，收敛版 DNA，对齐 ChatGPT 建议 70/20/10）
BIZ_PRIMARY   = RGBColor(0x33, 0x4B, 0x5E)   # 主色 深蓝灰（标题 / 正文强调 / 卡片结论）
BIZ_TEAL_N    = RGBColor(0x00, 0xA6, 0xC8)   # 辅助 科技青（章节标签 / 关键词）
BIZ_GOLD_N    = RGBColor(0xC9, 0x9A, 0x1A)   # 强调 金（编号 / 结论关键词，克制使用）
BIZ_BODY      = RGBColor(0x3F, 0x45, 0x4B)   # 正文 深灰
BIZ_AUX       = RGBColor(0x70, 0x77, 0x7E)   # 辅助文字 灰（副线 / 页码）
BIZ_CARD_LINE = RGBColor(0xE5, 0xE9, 0xED)   # 内容卡极淡边框（贴白底，1px 感）
CARD_BG_BIZ   = RGBColor(0xF2, 0xF2, 0xF2)   # 浅灰卡底（备用，未用于 card_grid）
DARK_TEXT     = RGBColor(0x33, 0x33, 0x33)   # 旧正文色（compare/profile 等沿用，未删）
MED_GRAY      = RGBColor(0x66, 0x66, 0x66)   # 旧辅助灰（compare/profile 等沿用，未删）
ACCENT        = BIZ_TEAL_N                   # 封面副标题强调青
PAGE_W = 13.333

def _center_x(w):
    return (PAGE_W - w) / 2.0

# ═══════════════════════════════════════════
# 画布 / 页管理（空白画布，不克隆任何模板）
# ═══════════════════════════════════════════
def new_deck(template=None, total=30):
    """创建空白 16:9 画布（白底深字，不依赖企业模板）。"""
    return blank_deck()

def first_slide(prs):
    """新增空白页作封面载体（不复用模板占位符）。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, WHITE)   # 保险：确保白底
    return s

def add_content(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

# ═══════════════════════════════════════════
# 版式原语（绝对坐标手绘，白底深字）
# ═══════════════════════════════════════════
def cover(slide, title, subtitle, tags=None, total=1, n=1):
    """封面：标题置顶左（非居中巨标题），品牌青短规则线，副标题强调青，深字白底。"""
    title_pt = fit_cover_title(title, 40, min_pt=24)   # 占屏≤72%
    tf = cover_line(slide, 0.8, 2.2, 11.73, title, title_pt, BIZ_PRIMARY, True, PP_ALIGN.LEFT)
    bar_y = shape_bottom_in(tf) + 0.14
    rect(slide, 0.8, bar_y, 1.6, 0.06, fill=BIZ_TEAL_N, line=None)
    tf = cover_line(slide, 0.8, bar_y + 0.06 + 0.20, 11.73, subtitle, 22, ACCENT, True, PP_ALIGN.LEFT)
    if tags:
        cover_line(slide, 0.8, shape_bottom_in(tf) + 0.16, 11.73,
                   "   ·   ".join(tags), 14, BIZ_AUX, False, PP_ALIGN.LEFT)
    cover_qr(slide, "business")
    add_page_number(slide, n, total, BIZ_AUX)

def section_header(slide, h1, label, sub, n=None, total=None):
    """行动标题：置顶左、结论式、22pt 深蓝粗体（McKinsey 标志）。"""
    tf = tb_box(slide, 0.46, 0.28, 11.5, 0.6)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    run_text(tf, h1, 22, True, BIZ_PRIMARY, PP_ALIGN.LEFT, Pt(0))
    # 副线：品牌青加粗标签 + 中灰说明（同一段落）；强制 MIDDLE 根治"框没包住/不居中"
    tf = tb_box(slide, 0.46, 0.95, 12.0, 0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = run_text(tf, label, 14, True, BIZ_TEAL_N, PP_ALIGN.LEFT, Pt(0), first=True)
    r = p.add_run(); r.text = "  " + sub
    r.font.size = Pt(14); r.font.bold = False; r.font.name = FONT
    r.font.color.rgb = BIZ_AUX; set_chinese_font(r)
    if n and total:
        add_page_number(slide, n, total, BIZ_AUX)

def timeline(slide, items, n=None, total=None, gold=None, gold_kw_a=()):
    """垂直时间轴：公共骨架 layout_timeline（竖线+圆点，皮肤 BUSINESS_SKIN.timeline）。"""
    layout_timeline(slide, BUSINESS_SKIN,
                    [{"date": it["time"], "text": it["text"],
                      "kw_teal": it.get("kw_teal", []), "kw_amber": it.get("kw_amber", []),
                      "amt": it.get("amt")}
                     for it in items],
                    y0=1.70, y_end=6.10)
    if gold:
        bottom_gold(slide, gold, n=n, total=total, kw_a=gold_kw_a)
    elif n and total:
        add_page_number(slide, n, total, MED_GRAY)
def two_column(slide, left, right, n=None, total=None, h=4.1, gold=None):
    """双栏：公共骨架 layout_two_column（透明底+青边框）。"""
    ye = 6.10 if gold else 6.30
    layout_two_column(slide, BUSINESS_SKIN, {"columns": [left, right]}, y0=1.70, y_end=ye)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, MED_GRAY)
def card_grid(slide, cards, cols=2, n=None, total=None):
    """卡片网格：公共骨架 layout_card_grid（透明底+青边框）。"""
    layout_card_grid(slide, BUSINESS_SKIN, cards, cols=cols, y0=1.50, y_end=6.38)
    if n and total:
        add_page_number(slide, n, total, MED_GRAY)
def summary(slide, title, sub, metrics=None, quote=None, n=None, total=None):
    """总结页：公共骨架 layout_summary（居中大标题+指标+金色收尾金句）。"""
    layout_summary(slide, BUSINESS_SKIN, {"title": title, "sub": sub,
                                          "metrics": metrics or [], "quote": quote},
                   y0=1.70, y_end=6.10)
    if n and total:
        add_page_number(slide, n, total, MED_GRAY)
def bottom_gold(slide, text, n=None, total=None, kw_a=()):
    """底部金句：金色居中收尾语（无分隔线）。机制委托 common.bottom_gold（共享金句坐标）。"""
    _gold(slide, text, color=BIZ_GOLD_N, size=16,
          page_color=BIZ_AUX, n=n, total=total, kw_a=kw_a)

# ═══════════════════════════════════════════
# 富内容原语（委托 common.py 内容专用骨架 + BUSINESS_SKIN；签名对齐 enterprise，
# 自含 section_header / 底部金句，runner 只需传内容。白底深字 McKinsey 调性。）
# ═════════════════════════════════════════
def definition_compare(slide, h1, label, hero, hero_en, hero_cn,
                       left_bullets, left_summary, compare_header,
                       negatives, positive, gold=None, n=None, total=None):
    """定义+排除对比：公共骨架 layout_compare（白底：深 hero / 青强调 / 浅灰圆底×badge）。"""
    section_header(slide, h1, label, "")
    layout_compare(slide, BUSINESS_SKIN, {
        "hero": hero, "hero_en": hero_en, "hero_cn": hero_cn,
        "bullets": left_bullets, "summary": left_summary,
        "right": {"header": compare_header,
                  "negatives": [{"kw": kw, "desc": d} for kw, d in negatives],
                  "positive": {"kw": "", "desc": positive}},
    }, y0=1.85, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, BIZ_AUX)

def work_people(slide, h1, label, loop_nodes, roles, abilities,
                gold=None, n=None, total=None):
    """「如何工作，谁适合做？」：公共骨架 layout_loop_page（①闭环 ②双卡⇄ ③能力卡）。"""
    section_header(slide, h1, label, "")
    content = {
        "cap1": "① 工作方式：业务现场的一体化闭环",
        "nodes": [nd if isinstance(nd, dict) else {"name": nd} for nd in loop_nodes],
        "cap2": "② 两类角色：边听边做，边做边改",
        "cards": [{"title": it["title"], "desc": it["desc"]} for it in roles],
        "cap3": "③ 三项核心能力",
        "tail": [{"title": it["title"], "desc": it["desc"]} for it in abilities],
    }
    layout_loop_page(slide, BUSINESS_SKIN, content, y0=1.88, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, BIZ_AUX)

def profile_warning(slide, h1, label, left, right, gold=None, n=None, total=None):
    """「人才画像 + 清醒提醒」：公共骨架 layout_profile_warning（非对称双卡+结论）。"""
    section_header(slide, h1, label, "")
    layout_profile_warning(slide, BUSINESS_SKIN, {"left": left, "right": right},
                           y0=1.85, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, BIZ_AUX)

def transform_summary(slide, h1, label, rows, gold=None, n=None, total=None):
    """「转型对照」：公共骨架 layout_transition_rows（旧态→新态，行高自适应）。"""
    section_header(slide, h1, label, "")
    layout_transition_rows(slide, BUSINESS_SKIN,
                           {"cap": "三个转变",
                            "rows": [{"from": it["from"], "to_title": it["to_title"],
                                      "to_lines": [it["to_desc"]]} for it in rows]},
                           y0=1.77, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, BIZ_AUX)

# ═══════════════════════════════════════════
# 防损清理（去除主题阴影 / 3D，保持矢量轻量）
# ═══════════════════════════════════════════
def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

def save(prs, outpath):
    """保存并防损清理。"""
    prs.save(outpath)
    full_cleanup(outpath)
    print("saved:", outpath, "slides:", len(prs.slides._sldIdLst))


# ═══════════════════════════════════════════
# 商务风皮肤（布局骨架的视觉参数，2026-07-31 架构：布局在 common.layout_*，皮肤在本层）
# ═══════════════════════════════════════════
BUSINESS_SKIN = {
    "bg": WHITE,
    "fonts": {"head": "微软雅黑", "body": "微软雅黑", "mute": BIZ_AUX},
    "cap":  {"size": 14, "bold": True, "color": BIZ_PRIMARY},
    "node": {"size": 16, "bold": True, "color": BIZ_PRIMARY, "sub_size": 12, "sub_color": BIZ_AUX},
    "card": {"title_size": 16, "title_color": BIZ_PRIMARY, "sub_size": 12, "sub_color": BIZ_AUX,
             "body_size": 12, "body_color": BIZ_BODY, "concl_color": BIZ_PRIMARY,
             "num_color": BIZ_GOLD_N,
             "line": BIZ_CARD_LINE, "line_w": 1.0,
             "shadow_alpha": 8000, "shadow_blur": 110000, "shadow_dist": 23000},
    "arrow": {"color": MED_GRAY, "w": Pt(1.5)},  # 白底 → 深色细线
    "timeline": {"axis": BIZ_TEAL, "dot": BIZ_TEAL, "date_size": 13, "date_color": BIZ_DARK,
                 "event_size": 14, "event_color": DARK_TEXT, "arrow": False,
                 "kw_teal_color": BIZ_TEAL, "kw_amber_color": GOLD},
    # 定义对比（白底：深 hero / 青强调 / 浅灰圆底×）
    "compare": {
        "hero_size": 48, "hero_color": BIZ_DARK,
        "en_size": 16, "en_color": MED_GRAY, "en_italic": True,
        "cn_size": 18, "cn_color": BIZ_DARK,
        "bullet_size": 14, "bullet_color": DARK_TEXT, "bullet_mark": "dot", "mark_color": BIZ_TEAL,
        "summary_size": 15, "summary_color": GOLD,
        "header_size": 18, "header_color": BIZ_DARK,
        "neg_badge": "circle", "neg_kw_size": 18, "neg_kw_color": BIZ_DARK,
        "neg_desc_size": 15, "neg_desc_color": MED_GRAY,
        "neg_badge_bg": RGBColor(0xEE, 0xF1, 0xF2), "neg_badge_color": MED_GRAY,
        "pos_bg": RGBColor(0xF2, 0xF8, 0xF9), "pos_border": BIZ_CARD_LINE,
        "pos_kw_size": 14, "pos_kw_color": BIZ_TEAL,
        "pos_desc_size": 14, "pos_desc_color": DARK_TEXT,
        "div": BIZ_TEAL, "right_card": False,
    },
    # 价值网格（白底：深字 / 金结论 / 等长循环箭头）
    "value_grid": {
        "num_size": 28, "num_color": BIZ_TEAL, "title_size": 16, "title_color": BIZ_DARK,
        "body_size": 12, "body_color": DARK_TEXT, "tag_size": 16, "tag_color": GOLD,
        "tag_align": "center", "div": MED_GRAY, "d1": False, "mark": False,
        "gap_h": 0.33, "gap_v": 0.16, "solve": False,
        "line": BIZ_CARD_LINE, "line_w": 1.25,
    },
    # 画像（浅灰实心卡 + 青标题线 + 金结论）
    "profile": {
        "card_bg": RGBColor(0xF7, 0xF9, 0xFA), "title_size": 18, "title_color": BIZ_DARK,
        "title_line": BIZ_TEAL, "role_size": 13, "role_color": MED_GRAY,
        "dot": BIZ_TEAL, "col_threshold": 8,
        "warn_kw_size": 16, "warn_kw_color": BIZ_DARK,
        "warn_desc_size": 12, "warn_desc_color": DARK_TEXT, "mark": BIZ_TEAL,
        "faint_div": RGBColor(0xE8, 0xEC, 0xEE), "strong_div": RGBColor(0xC9, 0xD1, 0xD6),
        "concl_bar": GOLD, "concl_size": 14, "concl_color": GOLD,
    },
    # 转型对照（旧态浅灰降温 / 新态青边高亮）
    "transition": {
        "from_bg": RGBColor(0xEE, 0xF1, 0xF2), "from_color": MED_GRAY,
        "to_color": BIZ_DARK, "line": BIZ_CARD_LINE,
    },
    "tail_bar": MED_GRAY,
    "line_spacing": 1.2,
}
