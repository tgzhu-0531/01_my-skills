# -*- coding: utf-8 -*-
"""
common.py — 四风格共享基件
提供：路径解析、品牌色、间距规则（已沉淀）、文字框/运行文本工具、二维码、页码。
所有风格引擎（enterprise / fengmang / wenzhi / business）均从此导入，保证间距与字体统一。
"""
import os
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Length
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ═══════════════════════════════════════════
# 路径解析（engines/common.py → skill 根目录）
# ═══════════════════════════════════════════
SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

def template_path():
    return os.path.join(SKILL_DIR, "templates", "enterprise.pptx")

def qr_path():
    return os.path.join(SKILL_DIR, "assets", "qrcode.jpg")

# ═══════════════════════════════════════════
# 品牌色（企业风基准，其他风格可覆盖）
# ═══════════════════════════════════════════
BIZ_TEAL    = RGBColor(0x00, 0xA7, 0xCB)   # 品牌青（主强调）
BIZ_TEAL2   = RGBColor(0x00, 0xAF, 0xD2)   # 主题强调青
BIZ_DARK    = RGBColor(0x38, 0x4F, 0x63)   # 深蓝强调
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE5, 0xE5, 0xE5)   # 正文浅灰
ORANGE      = RGBColor(0xF7, 0x96, 0x46)   # 警示橙
GOLD        = RGBColor(0xC9, 0xA2, 0x27)   # 金色（金句/高亮）
BLACK       = RGBColor(0x00, 0x00, 0x00)

# ═══════════════════════════════════════════
# 间距规则（已与用户确认沉淀，全风格统一）
# ═══════════════════════════════════════════
SP_UNIT      = Pt(4)    # 单位内：同主题连续条目之间的段后距
SP_GROUP     = Pt(24)   # 单位间：不同主题/模块之间的段后距
SP_CARD_BODY = Pt(14)   # 卡片标题 → 卡片内容
SP_RUN_DEF   = Pt(4)    # run_text 默认段后距
SP_SUBITEM   = Pt(6)    # 子项内部 标题→正文

# ── 画布与安全边距（四风格共用 · 强制统一；消除 0.46/0.60/0.80 漂移）──
CANVAS_X0   = 0.80     # 左安全边距
CANVAS_X1   = 12.53    # 右边界 = 13.333 - 0.80
CANVAS_CW   = 11.73    # 内容区宽 = 12.53 - 0.80
CANVAS_TOP  = 1.88     # 内容区上边界（标题 / 副标之下）
CANVAS_BOT  = 6.20     # 内容区下边界（页码之上）
CANVAS_CH   = CANVAS_BOT - CANVAS_TOP   # 内容区高 = 4.32
PAGE_W     = 13.333   # 整屏宽（16:9）；封面标题占屏比例基准
COVER_TITLE_MAX_W_RATIO = 0.72   # 封面主标题文字宽 ≤ 整屏宽×此比例（落实 SKILL §4.2.1）
# 底部金句文本统一坐标（由页底反推，四风格共享，杜绝散落硬编码 6.20/6.35/6.38/6.50 不一致）
# 基准取自企业风已验收值（text_y=6.65，距页底 0.85″，用户认定最佳）；其余三风原 6.38/6.50 偏高→统一下移到此
PAGE_H      = 7.5      # 整屏高（16:9）
GOLD_LINE_Y = 6.50     # 【已弃用 2026-08-01】金句分隔线 y——金句上方无分隔线（通用规则），仅保留兼容签名
GOLD_TEXT_Y = 6.65     # 金句文本 y（框高 0.6 → 框底 7.25，距页底 0.25″；与企业风观感对齐）

# ── 段落间距节奏（通用，替换原先散落的固定值）──
SP_LINE       = 1.22    # 段内行距倍数（leading ratio；替换 run_text 写死的 1.2）
NATURAL_LH    = 1.35    # 微软雅黑自然行高比率（PIL 实测 ascent+descent ≈1.35em）；
#                          测高地板：框高不得小于字形自然高，否则大字底部笔画被压（"框没包住文字"）
SP_PARA_INTRA = Pt(4)   # 段内多行额外上下留白
SP_PARA_INTER = Pt(24)  # 段落之间默认段后距（可被 auto_gap 覆盖）
SP_GAP_MIN    = Pt(10)  # 自动间距下限（防拥挤）
SP_GAP_MAX    = Pt(40)  # 自动间距上限（防过大白 / 散）

FONT = "微软雅黑"

# ═══════════════════════════════════════════
# 通用角色排版（跨风格默认层级；引擎可按版面覆盖单个角色）
# 字号单位 pt；加粗 bool。颜色不在此定义（见各引擎 PALETTE）。
# ═══════════════════════════════════════════
TYPOGRAPHY = {
    "hero":       (36, True),    # 封面主标题
    "h1":         (28, True),    # 页面/章节主标题
    "title_card": (18, True),    # 卡片标题
    "label":      (18, True),    # 副标题 / 核心观点
    "gold":       (16, True),    # 底部金句 / 结论（四风格统一 16pt，2026-08-01 由 20 收敛）
    "date":       (20, True),    # 时间轴日期
    "body":       (14, False),   # 正文（宽松区）
    "body_small": (12, False),   # 密排卡正文
    "badge":      (12, True),    # 徽章
}

# ═══════════════════════════════════════════
# 中文字体修复（关键：latin 只设英文，中文需显式设 ea/cs）
# ═══════════════════════════════════════════
def set_chinese_font(run, font=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', font)

# ═══════════════════════════════════════════
# 文字框 / 运行文本
# ═══════════════════════════════════════════
def _to_emu(v):
    """坐标单位自适应：已是 Length(Pt/Emu/Inches) 直接返回，否则按英寸转换。
    防止调用方传 Pt()(内部已是 EMU) 又被 Inches() 二次放大导致巨型形状。"""
    if isinstance(v, Length):
        return v
    return Inches(v)

def tb_box(slide, x, y, w, h):
    """创建文字框，margin 清零 + anchor=TOP + auto_size=NONE，返回 text_frame。"""
    box = slide.shapes.add_textbox(_to_emu(x), _to_emu(y), _to_emu(w), _to_emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf

def run_text(tf, text, size=14, bold=False, color=LIGHT_GRAY,
             align=PP_ALIGN.LEFT, space_after=SP_RUN_DEF, first=False):
    """在 text_frame 追加一段文本。first=True 且首段为空时复用首段。
    自动设 line_spacing=1.2（通用字号规范，专业 PPT 标准间距）。"""
    if first and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    p.space_before = Pt(0)
    p.line_spacing = SP_LINE
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    set_chinese_font(r)
    return p

def apply_role(tf, role, text, palette, align=PP_ALIGN.LEFT,
               space_after=SP_RUN_DEF, first=False):
    """按角色渲染文本：字号/加粗取自 TYPOGRAPHY，颜色取自 palette[role]。
    自动套中文（微软雅黑）字体修复。供各风格引擎统一调用。"""
    size, bold = TYPOGRAPHY[role]
    color = palette[role]
    return run_text(tf, text, size, bold, color, align, space_after, first)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    """画矩形（卡片底/标题线/装饰块）。坐标单位自适应（见 _to_emu）。"""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        _to_emu(x), _to_emu(y), _to_emu(w), _to_emu(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp

def rounded_card(slide, x, y, w, h, fill, adj=0.04, line=None, line_w=1.25):
    """实心圆角卡片（风格无关机制）。
    ⚠️ 关键坑：python-pptx ROUNDED_RECTANGLE 默认 adjustments[0]=0.16667，
    圆角太软（移动端 UI 感）；商务/企业等专业场景应设 0.03~0.05（克制小圆角）。
    fill: 卡底色(RGBColor，由引擎按风格传入)；line: 边框色(None=无边框)。
    自动去阴影 + 去 <p:style>（防 PowerPoint 主题样式干扰）。"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        _to_emu(x), _to_emu(y), _to_emu(w), _to_emu(h))
    card.adjustments[0] = adj
    card.shadow.inherit = False
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if line is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = line
        card.line.width = Pt(line_w)
    sp = card._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)
    _soft_shadow_if_light(card, slide)
    return card

# ═══════════════════════════════════════════
# 内容卡通用机制（带框卡类 · 跨风格通用，几何层，不绑配色）
# 设计原则见 design-principles.md §4.2：框体一律透明底 + 风格强调色亮边框；
# 框内四段式（标题→D1→正文→D2→结论），D1 可选、D2 仅当有结论。
# ═══════════════════════════════════════════
CARD_BORDER_W = Pt(1.25)   # 带框卡亮边框线宽（跨风格通用约定）

def content_card(slide, x, y, w, h, accent, line_w=1.25, adj=0.04,
                 shadow_alpha=8000, shadow_blur=90000, shadow_dist=20000):
    """带框内容卡（跨风格通用机制，几何层）：透明底 fill.background() + 风格极淡边框。
    线宽 1.25pt、小圆角 adj=0.04（克制不软，见 rounded_card 机制坑）。
    边框色 accent 由引擎传入本风格「极淡贴背景」色（不再用亮青）；浅底风格自动叠柔和投影(暗底跳过)。
    **仅此一种填充：无浅底例外、无实心底例外**（实心容器卡不属此类，见 §4.3）。
    accent 不写死。自动去 <p:style> 防主题样式干扰。
    shadow_alpha 默认 8000（8%，克制轻影；旧默认 35000=35% 偏重已弃用），个别风格可再调淡(如文质 6000)。
    详见 references/design-principles.md §11 卡片克制规范。"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        _to_emu(x), _to_emu(y), _to_emu(w), _to_emu(h))
    card.adjustments[0] = adj
    card.shadow.inherit = False
    card.fill.background()
    card.line.color.rgb = accent
    card.line.width = Pt(line_w)
    _clean_shape(card)
    _soft_shadow_if_light(card, slide, blur=shadow_blur, dist=shadow_dist, alpha=shadow_alpha)
    return card

def add_soft_shadow(shape, blur=90000, dist=20000, direction=5400000,
                    color='000000', alpha=8000):
    """给形状加柔和投影（标准 <a:outerShdw>），模拟卡片悬浮。
    仅写入 slide 级形状 XML，不会被 business 的 full_cleanup 剥掉（其只清 theme）。
    blur/dist 单位 EMU；direction 单位 1/60000 度（5400000=正下）；alpha 单位 1/1000 百分比。
    alpha 默认 8000（8%，克制轻影；旧默认 35000=35% 偏重已弃用）。"""
    spPr = shape._element.spPr
    for old in spPr.findall(qn('a:effectLst')):
        spPr.remove(old)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    outer = spPr.makeelement(qn('a:outerShdw'), {
        'blurRad': str(blur), 'dist': str(dist),
        'dir': str(direction), 'rotWithShape': '0'})
    srgb = spPr.makeelement(qn('a:srgbClr'), {'val': color})
    srgb.append(spPr.makeelement(qn('a:alpha'), {'val': str(alpha)}))
    outer.append(srgb)
    effectLst.append(outer)
    spPr.append(effectLst)
    return shape

def _soft_shadow_if_light(card, slide, blur=90000, dist=20000, alpha=8000):
    """浅底风格给卡片加柔和投影；暗底(如锋芒近黑)跳过，靠卡片填充差抬升。
    blur/dist/alpha 可传参（默认 8000=8% 轻影；旧 35000=35% 偏重已弃用）。
    暗底风格（亮度<90）一律不叠阴影，靠填充差抬升——避免深色阴影在深色底上不可见却增重。"""
    try:
        rgb = slide.background.fill.fore_color.rgb
        if 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2] < 90:
            return
    except Exception:
        pass
    add_soft_shadow(card, blur=blur, dist=dist, alpha=alpha)

def card_divider(slide, x, y, w, color, pt=0.75, full=True, title_w=None):
    """卡内分隔线（跨风格通用机制，几何层）。(x, y) 为线条中心。
    full=True  → 通栏（D2 结论上分界，宽=w，风格主分隔线色）；
    full=False → 不通栏（D1 标题下头身界，宽=title_w 或 w，强调色）。
    仅画一条横细线，颜色/线宽由引擎传入；pt 单位为点，内部转英寸。"""
    dw = w if full else (title_w if title_w is not None else w)
    th = pt / 72.0
    return rect(slide, x, y - th / 2.0, dw, th, fill=color, line=None)

# ═══════════════════════════════════════════
# 通用布局约束（四风格共用 · 强制）
# 设计原则见 design-principles.md §10：
#  ① 画布安全边距统一（消除 0.46/0.60/0.80 漂移）
#  ② 正文居页中（内容区垂直居中，杜绝底部大白）
#  ③ 间距按内容自动计算（少→放大填满不空，多→收紧防挤）
#  ④ 段内行距 / 段间间距统一节奏（SP_*）
# 仅几何，不绑配色 / 字体，不影响各风格视觉身份。
# ═══════════════════════════════════════════
def vcenter(total_h, top=CANVAS_TOP, bot=CANVAS_BOT):
    """内容块在 [top,bot] 内垂直居中起始 y。total_h 为内容块总高(英寸)。
    内容超出可用区时回退到 top（宁可紧凑也不鬼畜居中）。"""
    avail = bot - top
    if total_h >= avail:
        return top
    return top + (avail - total_h) / 2.0

def _vshift(slide, n0, y0, yb, y_end):
    """将本次调用新增的形状（slide.shapes[n0:]）整体在 [y0,y_end] 内垂直居中。
    内容块实际占 [y0, yb]；若 yb<y_end，则下移 delta 使其居中。仅移新增形状，不动页眉/页码。"""
    block_h = yb - y0
    avail = y_end - y0
    if block_h <= 0 or block_h >= avail:
        return
    delta = (avail - block_h) / 2.0
    # 注意：python-pptx 对 ShapeTree 做切片 [n0:] 会返回裸 lxml 元素而非 Shape 对象，
    # 必须用 list() 强制物化后才能取到带 .top 的 Shape 对象。
    all_shapes = list(slide.shapes)
    for sh in all_shapes[n0:]:
        if sh.top is not None:
            sh.top = Inches(sh.top / 914400.0 + delta)

def auto_gap(total_blocks_h, n_blocks, top=CANVAS_TOP, bot=CANVAS_BOT,
             gmin=SP_GAP_MIN, gmax=SP_GAP_MAX):
    """按内容量自动算块 / 段间距(英寸)：内容少→间距放大填满(不留大白)，
    内容多→间距收紧(不拥挤)。夹在 [gmin,gmax] 内防极端。n_blocks<=1 返回 0。"""
    if n_blocks <= 1:
        return 0.0
    avail = bot - top
    slack = avail - total_blocks_h
    if slack <= 0:
        return gmin.pt
    gap = slack / (n_blocks - 1)
    return max(gmin.pt, min(gmax.pt, gap))

def set_leading(tf, ratio=SP_LINE):
    """统一段内行距（替换各引擎散落的 0.93 / 1.35 / 1.2 写死值，防漂移）。"""
    for p in tf.paragraphs:
        p.line_spacing = ratio

def measure_para_h(text, pt, width, leading=SP_LINE):
    """估算一段文本在给定框宽内的渲染高度(英寸)：按实测字宽折行 × 行距。"""
    if not text or not text.strip():
        return pt * leading / 72.0
    tw = measure_text_width(text, pt)
    n_lines = max(1, math.ceil(tw / max(width, 0.1)))
    return n_lines * pt * leading / 72.0

def vstack(slide, blocks, x, w, top=CANVAS_TOP, bot=CANVAS_BOT,
           gap_mode='auto', gmin=SP_GAP_MIN, gmax=SP_GAP_MAX):
    """竖向自动布局（几何层）。blocks=[(height_in, draw_fn(y)->None), ...]；
    整体在 [top,bot] 内垂直居中，块间用 auto_gap 自动间距（填满不空、夹紧不挤）。
    draw_fn 内用本风格配色 / 字体，本函数只定几何。"""
    n = len(blocks)
    if n == 0:
        return
    total = sum(h for h, _ in blocks)
    if gap_mode == 'auto':
        gap = auto_gap(total, n, top, bot, gmin, gmax)
    else:
        gap = SP_PARA_INTER.pt
    y = vcenter(total + (n - 1) * gap, top, bot)
    for h, fn in blocks:
        fn(y)
        y += h + gap

def fit_content_font(texts, box_w_in, box_h_in, font_name=FONT, bold=False,
                     line_spacing=SP_LINE, gap_in=0.04, min_size=12, max_size=20,
                     max_lines=None):
    """通用规则 §10.6：内容区字号按版面自适配，下限 min_size(默认12)。

    主标题 / 副标题 / 金句字号固定（各引擎角色表），**内容区**文本（正文/卡片标题/要点等）
    可调用本函数：按给定框（宽 box_w_in × 高 box_h_in 英寸）内实测字宽折行，
    取最大的整数 pt ∈ [min_size, max_size] 使总渲染高 ≤ 框高。
    内容稀疏→字号放大更大气；内容密集→收紧但不低于下限。
    max_lines：单段折行上限（如 2 = 每段不允许超过 2 行），超出则该字号不合格——
    与总高约束正交：既防"一堵字墙"（行数上限），又防"底部挤压"（总高上限）。
    texts 为单条字符串或字符串列表（列表按整组计算，组内间距 gap_in 计入；
    同一 2×2 网格等并列卡建议取整组最小值保证同级一致）。
    line_spacing 须与目标渲染器的实际行距一致（如 wenzhi block 用 1.2），
    font_name 由各风格传入，不绑死；无字体时 measure_text_width 兜底估算。"""
    if isinstance(texts, str):
        texts = [texts]
    def wrapped_lines(text, S):
        if not text.strip():
            return 1
        return max(1, math.ceil(measure_text_width(text, S) * 1.03 / max(box_w_in, 0.1)))
    def fits(S):
        ls = [wrapped_lines(t, S) for t in texts]
        if max_lines is not None and any(l > max_lines for l in ls):
            return False
        th = sum(ls) * S * line_spacing / 72.0 + (len(texts) - 1) * gap_in
        return th <= box_h_in
    best = min_size
    for S in range(min_size, max_size + 1):
        if fits(S):
            best = S
        else:
            break
    return best

# ═══════════════════════════════════════════
# 二维码 / 页码 / 画布
# ═══════════════════════════════════════════
def add_qr(slide, x, y, size=1.1, path=None):
    path = path or qr_path()
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(size), Inches(size))

def add_page_number(slide, n, total, color=LIGHT_GRAY, x=12.2, y=7.1, size=9):
    """页码：9pt 浅色右对齐，y=7.1（距底 0.4in）。"""
    tf = tb_box(slide, x, y, 1.0, 0.4)
    run_text(tf, f"{n} / {total}", size, False, color, PP_ALIGN.RIGHT, Pt(0), first=True)

# ═══════════════════════════════════════════
# 文字框高度自适应（落实 SKILL §4.12-6「尺寸由内容实测」对高度的要求）
# 与 measure_text_width 成对：宽度实测算，高度也实测算，框高紧贴文字+小间隙
# ═══════════════════════════════════════════
def measure_text_height(text, pt, width_in=None, line_spacing=SP_LINE, font_path=None):
    """测量文本渲染高度（英寸）。
    复用 measure_text_width 估算折行：给定框宽则按实测最长行宽推算行数，否则单行。
    pt 字号；line_spacing 行距倍数；返回英寸。供 tb_box_hug 让框高紧贴文字。
    ⚠️ 每行高取 max(line_spacing, NATURAL_LH)：PowerPoint 不会把行渲染得比字形自然高更窄，
       行距倍数偏小时（如 1.22 < 1.35）会被字形自然行高顶起，测高必须用地板，否则大字'框没包住'。"""
    if not text:
        return 0.0
    if width_in and width_in > 0:
        # 每行可容纳像素宽 ≈ 框宽*0.96（预留边距）；按实测字宽推算行数
        capacity = width_in * 72.0 * 0.96
        total_px = measure_text_width(text, pt, font_path) * 72.0
        lines = max(1, math.ceil(total_px / capacity))
    else:
        lines = 1
    eff = max(line_spacing, NATURAL_LH)   # 自然行高地板，防大字底部笔画被压
    return lines * pt * eff / 72.0

def tb_box_hug(slide, x, y, w, text, pt, color=LIGHT_GRAY, bold=False,
               align=PP_ALIGN.LEFT, line_spacing=SP_LINE, gap=0.06,
               anchor=MSO_ANCHOR.TOP, first=False, space_after=Pt(0)):
    """创建「高度紧贴文字」的文字框：高度 = 实测文字高 + gap，保持 auto_size=NONE、margin 0。
    替代写死的偏大固定 h（落实 SKILL §4.4 框高紧贴）。单段落文本；多段落仍用 tb_box。返回 text_frame。"""
    h = measure_text_height(text, pt, w, line_spacing) + gap
    tf = tb_box(slide, x, y, w, h)
    tf.vertical_anchor = anchor
    run_text(tf, text, pt, bold, color, align, space_after, first)
    return tf

def shape_bottom_in(tf):
    """返回 tb_box/tb_box_hug 所属 shape 的底边 y（英寸）。供 flow 布局按前一元素实际底推导下一元素 y。"""
    sh = tf._parent
    return (sh.top + sh.height) / 914400.0

def fit_font_to_width(text, max_pt, min_pt=12, max_w=None, font_path=None):
    """按可用宽反解「能单行容纳的最大字号」：从 max_pt 逐级下探，直到实测宽 ≤ max_w。
    max_w 未给或文本为空则返回 max_pt（不缩）。返回选定字号(int)。
    根治封面/主标题固定字号导致长文本换行（落实 SKILL §4.2.1 首页主标题、§4.2.2 内容页主标题'单行·字号智能算'）。"""
    if not text or not max_w or max_w <= 0:
        return max_pt
    pt = int(max_pt)
    while pt > min_pt and measure_text_width(text, pt, font_path) > max_w:
        pt -= 1
    return max(pt, min_pt)

def cover_line(slide, x, y, w, text, pt, color=LIGHT_GRAY, bold=False,
               align=PP_ALIGN.LEFT, font=None, gap=0.06, line_spacing=SP_LINE,
               cy=None):
    """封面文本公共原语：框高 = 实测字高 + 2×gap（上下等量呼吸）+ 强制 vertical_anchor=MIDDLE。
    根治四风格封面『框没包住文字』或『文字不居中』不一致（落实 SKILL §4.2.1）。
    font：可传字体名（文质 hero 用宋体衬线，不能强套雅黑）；缺省=雅黑。
    cy：若给定，y 视为文字垂直中心（框顶 = cy - h/2），用于保持既有视觉重心不变；
        否则 y 为框顶，配合 shape_bottom_in 做 flow 布局。返回 text_frame。
    ⚠️ 与 tb_box_hug 区别：hug 用 anchor 透传(各家手写出错)、gap 仅底部；本原语强制居中+等量呼吸，统一兜底。"""
    h = measure_text_height(text, pt, w, line_spacing) + 2 * gap
    if cy is not None:
        y = cy - h / 2.0
    tf = tb_box(slide, x, y, w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, text, pt, bold, color, align, Pt(0), first=True)
    if font:
        set_chinese_font(tf.paragraphs[0].runs[0], font)
    return tf

def fit_cover_title(text, max_pt, min_pt=12, ratio=COVER_TITLE_MAX_W_RATIO,
                    page_w=PAGE_W, content_w=CANVAS_CW, font_path=None):
    """封面主标题字号反解：max_w = min(内容区宽, 页宽×比例)。
    根治『标题顶满整屏』（落实 SKILL §4.2.1 首页主标题占屏≤比例）。返回区间 [min_pt, max_pt] 内最大字号。"""
    max_w = min(content_w, page_w * ratio)
    return fit_font_to_width(text, max_pt, min_pt, max_w, font_path)

# 风格→封面二维码坐标（SKILL §4.9 坐标表；底中，与金句保留 0.3in 间隙）
COVER_QR = {
    "enterprise": (6.12, 5.55, 1.0),
    "fengmang":   (6.12, 5.85, 1.0),
    "wenzhi":     (6.12, 5.55, 1.05),
    "business":   (6.12, 5.55, 1.0),
}

def cover_qr(slide, style, label="扫码关注「天戈朱」"):
    """封面底中二维码 + 提示标签（SKILL §4.9）。style 取 COVER_QR 表坐标；无匹配则跳过。"""
    if style not in COVER_QR:
        return
    x, y, size = COVER_QR[style]
    add_qr(slide, x, y, size)
    tf = tb_box(slide, x - 0.3, y + size + 0.05, size + 0.6, 0.3)
    run_text(tf, label, 9, False, LIGHT_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)

def bottom_gold(slide, text, color, line_color=None, size=16, bold=True,
                x=0.8, w=11.73, line_y=GOLD_LINE_Y, text_y=GOLD_TEXT_Y, box_h=0.6,
                n=None, total=None, page_color=None, align=PP_ALIGN.CENTER,
                kw_t=(), kw_a=()):
    """底部金句（公用实现，去重四风格）。
    各风格传自己配色/字号/bold；page_color 缺省=LIGHT_GRAY。
    坐标默认走共享常量 GOLD_TEXT_Y（由页底反推，四风格一致，落实 SKILL §4.2.2）。
    ⚠️ 强制 vertical_anchor=MIDDLE：根治四风格金句'框没包住文字/文字不上下居中'。
    bold 默认 True：四风格金句统一加粗。
    kw_t/kw_a：金句内关键词高亮（如金额"75亿美元"→kw_a 琥珀加粗），落实内容高亮规则。
    2026-08-01 通用规则：页底金句上方无分隔线（line_color/line_y 参数弃用，仅保留签名兼容调用方）。"""
    _guard_one_line(text, size, w, ctx="页底金句")   # 硬规则：金句强制单行、禁换行
    tf = tb_box(slide, x, text_y, w, box_h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kw_t or kw_a:
        emph_runs(tf, text, size, color, kw_t, kw_a,
                   color_t=color if not kw_t else BIZ_TEAL,
                   color_a=GOLD, font=FONT)
    else:
        run_text(tf, text, size, bold, color, align, Pt(0), first=True)
    if n and total:
        add_page_number(slide, n, total, page_color or LIGHT_GRAY)


def emph_runs(tf, text, size, base_color, kw_t=(), kw_a=(),
              color_t=None, color_a=None, font=None):
    """正文关键词高亮公共助手（落实 SKILL §4 内容页高亮规则）：
    - 最多 2 种高亮色：kw_t→青(BIZ_TEAL)、kw_a→金(GOLD)；其余保持正文色
    - 克制：仅首次出现高亮（同词本段去重），滥则乱
    - 长词优先切分，避免子串提前截断
    返回 None；tf 须已建好（本助手只填 runs）。"""
    import re
    all_kw = list(kw_t) + list(kw_a)
    f = font or FONT
    if not all_kw:
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(size); r.font.bold = False
        r.font.name = f; r.font.color.rgb = base_color
        set_chinese_font(r, f)
        return
    all_kw.sort(key=len, reverse=True)
    pat = re.compile("(" + "|".join(re.escape(k) for k in all_kw) + ")")
    parts = [p for p in pat.split(text) if p]
    seen = set()
    ct = color_t or BIZ_TEAL
    ca = color_a or GOLD
    kw_t_set, kw_a_set = set(kw_t), set(kw_a)
    for part in parts:
        r = tf.paragraphs[0].add_run()
        r.text = part
        r.font.size = Pt(size); r.font.name = f
        set_chinese_font(r, f)
        if part not in seen:
            if part in kw_t_set:
                r.font.bold = True; r.font.color.rgb = ct
            elif part in kw_a_set:
                r.font.bold = True; r.font.color.rgb = ca
            else:
                r.font.bold = False; r.font.color.rgb = base_color
            seen.add(part)
        else:
            r.font.bold = False; r.font.color.rgb = base_color


def blank_deck():
    """创建 16:9 空白演示文稿（锋芒/文质/商务 从零绘制用）。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    """设置幻灯片背景纯色填充。"""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

# ═══════════════════════════════════════════
# 通用工具：占位符清理 / shape 清理
# ═══════════════════════════════════════════
def _clean_shape(shape):
    """移除 shape 的元素级 <p:style> 样式（防 PowerPoint 自带样式干扰）。"""
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

def remove_placeholders(slide):
    """XML 级别删除 slide 上全部占位符（彻底，不是 shape.text=''）。
    企业风封面必须调用——占位符用 text='' 清空后仍显示提示语。"""
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldXml = slide._element
    spTree = sldXml.find(qn('p:cSld')).find(qn('p:spTree'))
    to_remove = []
    for sp in spTree.findall('p:sp', nsmap):
        nvSpPr = sp.find('p:nvSpPr', nsmap)
        if nvSpPr is not None:
            nvPr = nvSpPr.find('p:nvPr', nsmap)
            if nvPr is not None and nvPr.find('p:ph', nsmap) is not None:
                to_remove.append(sp)
    for sp in to_remove:
        spTree.remove(sp)

# ═══════════════════════════════════════════
# 通用几何 / 工具（四风格引擎均可调用，不绑死企业风色板）
# ═══════════════════════════════════════════
def add_homeplate(slide, x, y, w, h, line_color=None, line_w=1.0, fill_color=None,
                  adj_pct=50, flip_h=False):
    """homePlate 预设几何（一侧切角五边形）。flip_h=True 切角在左，默认切角在右。
    用 OOXML 直接写 <a:prstGeom prst="homePlate">，避免旋转 PENTAGON 导致宽高被交换的坑。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    sp = shp._element
    spPr = sp.find(qn('p:spPr'))
    existing_geom = spPr.find(qn('a:prstGeom'))
    if existing_geom is not None:
        spPr.remove(existing_geom)
    prstGeom = etree.SubElement(spPr, qn('a:prstGeom'))
    prstGeom.set('prst', 'homePlate')
    avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    gd = etree.SubElement(avLst, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj_pct * 1000}')
    if flip_h:
        for xf in sp.iter(qn('a:xfrm')):
            xf.set('flipH', '1')
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    _clean_shape(shp)
    return shp

def add_dashed_line(slide, x, y, length, color, dash_w=0.5, horizontal=False):
    """短虚线（默认竖向）。用于时间轴每张卡片对齐的装饰线。
    (x,y) 起点，length 长度(in)。horizontal=True 时画横线。"""
    if horizontal:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x + length), Inches(y))
    else:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            Inches(x), Inches(y), Inches(x), Inches(y + length))
    conn.line.color.rgb = color
    conn.line.width = Pt(dash_w)
    dln = conn.line._get_or_add_ln()
    prstDash = etree.SubElement(dln, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    return conn

# ═══════════════════════════════════════════
# 闭环 / 箭头通用机制（风格无关；颜色由引擎传入，不写死）
# 线宽 2pt = 跨风格通用约定（几何规则，与底色无关）；
# 颜色须随风格底色走：深底风格(企业风深蓝 / 锋芒纯黑)用浅色 / 白色，
# 浅底风格(文质奶油纸 / 商务白底)用深色，否则白箭头在浅底直接隐形。
# ═══════════════════════════════════════════
LOOP_ARROW_W = Pt(2)   # 闭环 / 价值箭头线宽（跨风格通用）

def loop_arrow(slide, x1, y1, x2, y2, color=WHITE, w=LOOP_ARROW_W, arrow=True):
    """通用直线箭头 / 线段（风格无关）。color 由引擎传入（须与底色对比）；
    arrow=True 时末端加三角箭头。各引擎闭环 / 价值原语复用此机制，不烘焙风格颜色。"""
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = color
    cxn.line.width = w
    if arrow:
        ln = cxn.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return cxn

def loop_return_u(slide, centers, node_bot, u_bot, color=WHITE, w=LOOP_ARROW_W):
    """U 型闭环回流（跨风格通用机制，从企业风 work_people 提炼，2026-07-31 用户拍板公共化）。
    三段直线：起点节点底中 → 下 → 沿底部左 → 上进入终点节点底中（箭头指入）；
    起点/终点均节点底中、左右对称。centers — 各节点底中 x 列表（首尾为终点/起点）；
    node_bot — 节点底 y；u_bot — U 底部 y（节点下方留缝）。color/w 由引擎传入（随底色）。"""
    x4, x1 = centers[-1], centers[0]
    loop_arrow(slide, x4, node_bot, x4, u_bot, color=color, w=w, arrow=False)
    loop_arrow(slide, x4, u_bot, x1, u_bot, color=color, w=w, arrow=False)
    loop_arrow(slide, x1, u_bot, x1, node_bot, color=color, w=w, arrow=True)


def validate_deck(path, qr_on_cover=True, bg_colors=None, verbose=True):
    """通用校验（安全网，不约束视觉）。检查不变量：
    1) 每页无形状越界（≤ 13.333 × 7.502）
    2) 封面有图片(二维码) [qr_on_cover]
    3) 所有带箭头的 connector(a:tailEnd) 线宽 ≈ 2pt（跨风格通用约定）
    4) 若提供 bg_colors[页序] = RGBColor，则箭头色 ≠ 底色（防浅底误用白）
    返回 (ok, report_lines)。"""
    PAGE_W, PAGE_H = 13.333, 7.502
    p = Presentation(path)
    rep, ok = [], True
    # 1) 越界
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            r = (sh.left + sh.width) / 914400
            b = (sh.top + sh.height) / 914400
            if r > PAGE_W + 0.02 or b > PAGE_H + 0.02:
                ok = False
                rep.append(f"  [页{i}] 越界 右={r:.2f} 底={b:.2f}")
    # 2) 封面二维码
    if qr_on_cover and p.slides:
        pics = [sh for sh in p.slides[0].shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not pics:
            ok = False
            rep.append("  [封面] 无二维码图片")
    # 3)+4) 箭头
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if sh._element.tag.split('}')[-1] != 'cxnSp':
                continue
            ln = sh._element.find('.//' + qn('a:ln'))
            if ln is None or ln.find(qn('a:tailEnd')) is None:
                continue   # 仅检箭头（虚线装饰无 tailEnd，跳过）
            w = sh.line.width
            if w is not None:
                wpt = w.pt
                if not (1.5 <= wpt <= 2.5):
                    ok = False
                    rep.append(f"  [页{i}] 箭头线宽={wpt:.2f}pt (期望≈2)")
            if bg_colors and (i - 1) in bg_colors:
                try:
                    if sh.line.color.type == MSO_COLOR_TYPE.RGB:
                        hexv = '%02X%02X%02X' % tuple(sh.line.color.rgb)
                        bg = bg_colors[i - 1]
                        if hexv == '%02X%02X%02X' % tuple(bg):
                            ok = False
                            rep.append(f"  [页{i}] 箭头色=底色(隐形)")
                except Exception:
                    pass
    # 5) 单行文本框是否换行（副标题 label / 金句 gold 等高度≤0.6 且无显式换行）
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text
            if not txt.strip() or '\n' in txt:
                continue
            if sh.height is None or sh.width is None:
                continue
            if sh.height / 914400 > 0.6:
                continue   # 多行卡正文，非单行框，跳过
            fs = None
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        fs = run.font.size.pt if fs is None else max(fs, run.font.size.pt)
            if fs is None:
                fs = 18
            tw = measure_text_width(txt, fs)
            bw = sh.width / 914400
            if bw <= 0:
                continue
            # 估算文本在框宽内会折成的行数；若框高已足够容纳折行则不报
            # （避免误伤「窄栏多行」框——其高度本就按折行数测算，并无溢出）。
            n_est = max(1, math.ceil(tw / bw))
            line_h = fs * 1.2 / 72.0          # 与 run_text 的 line_spacing=1.2 一致
            req_h = n_est * line_h * 0.95     # 留 5% 容差
            act_h = sh.height / 914400
            if act_h < req_h - 0.03:
                ok = False
                rep.append(f"  [页{i}] 单行文本疑似换行 文本宽={tw:.2f} > 框宽={bw:.2f}  '{txt[:18]}...'")
    # 6) 卡片阴影透明度上限（克制轻影，非阻断 WARN）
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            outer = sh._element.find('.//' + qn('a:outerShdw'))
            if outer is None:
                continue
            alpha_el = outer.find('.//' + qn('a:alpha'))
            if alpha_el is None:
                continue
            try:
                av = int(alpha_el.get('val'))
            except Exception:
                continue
            if av > 12000:   # >12% 视为偏重（克制上限 8%）
                rep.append(f"  [页{i}] ⚠ 阴影偏重 alpha={av/1000:.0f}% (>12%)，建议 ≤8%（design-principles.md §11）")
    # 注：#6 为 WARN，不置 ok=False，避免阻断既有 deck；属设计纪律提醒。
    # 7) 内容底 ≤ 卡框底（抓文字溢出卡片，如 value_loop 结论溢出卡框）
    for i, s in enumerate(p.slides, 1):
        cards = [sh for sh in s.shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and sh.left is not None and sh.width is not None]
        for sh in s.shapes:
            if not sh.has_text_frame or sh.left is None or sh.width is None:
                continue
            if not sh.text_frame.text.strip():   # 空文本框（如满高装饰竖条）无文字可溢出，跳过
                continue
            cx = sh.left + sh.width / 2.0
            cy = sh.top + sh.height / 2.0
            tb = (sh.top + sh.height) / 914400.0
            best = None
            for c in cards:
                cl, ct, cw_, ch_ = c.left, c.top, c.width, c.height
                if cl <= cx <= cl + cw_ and ct <= cy <= ct + ch_:
                    area = (cw_ / 914400.0) * (ch_ / 914400.0)
                    if best is None or area < best[1]:
                        best = (c, area)
            if best is not None:
                card_b = (best[0].top + best[0].height) / 914400.0
                if tb > card_b + 0.08:
                    ok = False
                    rep.append(f"  [页{i}] 文字溢出卡框 底={tb:.2f} > 卡底={card_b:.2f}  '{sh.text_frame.text[:16]}'")
                elif tb > card_b + 0.03:
                    rep.append(f"  [页{i}] ⚠ 文字贴近卡框底 底={tb:.2f} 卡底={card_b:.2f}")
    # 8) label 分类词模式软告警（企业风 label 应为核心观点，勿用「中文·English」分类）
    #    仅当「·」前存在非空中文分类词时判定为分类词，避免误伤「· English 中文」型要点 bullets。
    for i, s in enumerate(p.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            if '·' in t:
                before, _, after = t.partition('·')
                before, after = before.strip(), after.strip()
                if before and after and after[0].isascii() and after[0].isalpha() and len(t) <= 28:
                    rep.append(f"  [页{i}] ⚠ label 疑似分类词(应为核心观点)  '{t}'")
    # 9) 垂直居中轻量启发（抓顶重/底重，仅告警）
    SLIDE_MID = 7.502 / 2.0
    for i, s in enumerate(p.slides, 1):
        if i == 1:
            continue
        tops = [sh.top / 914400.0 for sh in s.shapes
                if sh.has_text_frame and sh.top is not None and sh.text_frame.text.strip()]
        if len(tops) < 3:
            continue
        mid = (min(tops) + max(tops)) / 2.0
        if abs(mid - SLIDE_MID) > 1.2:
            rep.append(f"  [页{i}] ⚠ 内容块{'顶重' if mid < SLIDE_MID else '底重'}（中心 {mid:.2f} 偏离页心 {SLIDE_MID:.2f}）")
    # 10) 内容页品牌装饰存在性（抓静默丢失 logo/标题带；企业风克隆品牌版式，品牌在版式层）
    #     注意：add_slide(layout) 后 logo/标题带是「版式层形状」，不在 slide XML 里，
    #     故须检查 slide.slide_layout 而非 slide.shapes。软告警，不阻断生成。
    #     风格感知：仅当整份 deck 的版式/母版「存在」品牌形状时才逐页检查；
    #     商务风等纯代码绘制（空白版式、无模板品牌）整份无品牌 → 整体跳过，避免误报。
    def _has_brand(sh):
        try:
            _st = sh.shape_type
        except Exception:
            _st = None
        return _st == MSO_SHAPE_TYPE.PICTURE or (sh.name and '组合' in (sh.name or ''))
    _brand_anywhere = False
    for _m in p.slide_masters:
        _shapes = list(_m.shapes)
        for _lay in _m.slide_layouts:
            _shapes += list(_lay.shapes)
        if any(_has_brand(_sh) for _sh in _shapes):
            _brand_anywhere = True
            break
    if _brand_anywhere:
        for i, s in enumerate(p.slides, 1):
            if i == 1:
                continue
            lay = s.slide_layout
            if lay is None:
                rep.append(f"  [页{i}] ⚠ 无版式引用，无法核验品牌装饰")
                continue
            has_logo = any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in lay.shapes)
            has_band = any((sh.name and '组合' in sh.name) for sh in lay.shapes)
            if not has_logo:
                rep.append(f"  [页{i}] ⚠ 内容页版式缺 logo（品牌静默丢失？）")
            if not has_band:
                rep.append(f"  [页{i}] ⚠ 内容页版式缺顶部标题带（品牌静默丢失？）")
    if verbose:
        print("✅ validate_deck 通过" if ok else "❌ validate_deck 未通过")
        for line in rep:
            print(line)
    return ok, rep

def fmt_date_5digit(t):
    """日期统一 5 位：4/23 → 04/23（已 5 位则保持）。无法解析时原样返回。"""
    try:
        m, d = t.split('/')
        return f"{int(m):02d}/{int(d):02d}"
    except Exception:
        return t

# ═══════════════════════════════════════════
# 字体度量（可移植 / 跨机器：克隆后任意环境均可跑）
# 卡宽 / 文本框宽度由"最长文本行真实字体宽度"算出，不凭感觉、不撑满。
# 优先 PIL 真实度量；无字体或无 PIL 时用字符估算兜底，保证任何机器都能生成。
# ═══════════════════════════════════════════
def find_cjk_font():
    """跨 OS 探测中文字体路径；找不到返回 None（调用方用估算兜底）。"""
    candidates = [
        # Windows
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\msyhbd.ttc",
        r"C:\Windows\Fonts\simsun.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
        # Linux
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
    ]
    for f in candidates:
        if os.path.exists(f):
            return f
    return None

_CJK_FONT_CACHE = {}

def _load_pil_font(path, pt):
    """加载 PIL 字体（ttc 用 index 0）。失败返回 None。"""
    if path is None:
        return None
    key = (path, pt)
    if key in _CJK_FONT_CACHE:
        return _CJK_FONT_CACHE[key]
    font = None
    try:
        from PIL import ImageFont
        if path.lower().endswith(".ttc"):
            font = ImageFont.truetype(path, pt, index=0)
        else:
            font = ImageFont.truetype(path, pt)
    except Exception:
        font = None
    _CJK_FONT_CACHE[key] = font
    return font

def measure_text_width(text, pt, font_path=None):
    """测量文本渲染宽度（英寸）。
    优先用 PIL 真实字体度量（跨机器一致，前提是装了 PIL 且有 CJK 字体）；
    无字体 / 无 PIL 时按字符估算兜底（中文≈1.0*pt、ASCII≈0.55*pt，与微软雅黑实测接近），
    保证克隆到任意机器都能生成、不崩溃。
    pt 为字号；返回英寸（px/72）。"""
    if font_path is None:
        font_path = find_cjk_font()
    font = _load_pil_font(font_path, pt) if font_path else None
    if font is not None:
        try:
            px = font.getlength(text)
            return px / 72.0
        except Exception:
            pass
    # 兜底估算
    width_px = 0.0
    for ch in text:
        width_px += pt * (1.0 if ord(ch) > 0x2E80 else 0.55)
    return width_px / 72.0

def fit_single_line(tf, text, pt, buffer=0.2, max_w=None):
    """把单行文本框宽度设为「实测文本宽 + buffer」，根治框宽写死偏窄导致的换行
    （如副标题 label / 金句 gold 文本较长时，在框内折行而页面整体仍有余量）。
    通过 tf._parent 修改所属 shape 宽度；左缘由调用方负责（居中需另行设 left）。
    返回最终框宽(英寸)。max_w：框宽上限（如内容区可用宽）；超出则截断到上限
    （此时仍可能换行，由 validate_deck 第 5 项捕获）。"""
    w = measure_text_width(text, pt) + buffer
    if max_w is not None:
        w = min(w, max_w)
    try:
        sh = tf._parent            # python-pptx: TextFrame 的拥有 shape
        sh.width = Inches(w)
    except Exception:
        pass
    return w

def _guard_one_line(text, size, w, ctx=""):
    """硬规则：金句（页底 bottom_gold + 卡内结论 hl）必须单行，禁止换行。

    若实测需 >1 行，直接抛错——强制作者**精简文案到 1 行**，而非靠压字号补救
    （压字号会造成 13/14 这类隐形不一致，正是用户反对的）。
    落实用户「金句不允许换行」红线（SKILL §4.12-7）。"""
    rows = max(1, math.ceil(measure_text_width(text, size) * 1.04 / w))
    if rows > 1:
        raise ValueError(
            f"[金句禁换行] {ctx}：『{text}』在 {size}pt、可用宽 {w:.2f}″ 下需 {rows} 行，"
            f"请精简文案到 1 行（不要靠压字号）。")


# ═══════════════════════════════════════════
# 约束求解排版（跨风格通用机制，2026-07-31 用户拍板，价值网格迭代沉淀）
# 规则见 SKILL.md §4.12-7「布局规则（公共，强制）」：
#   · 上下间距按内容量智能算（区间内弹性：内容多→下限、内容少→放大，cap 各自上限）
#   · 间距层级：标题↔正文 ≥ 段落↔段落 ≥ 段内行距（由各间距区间 min 值体现）
#   · 正文区字号智能算：≥ 底线（默认12），从风格上限向下试，禁止手工指定
#   · 行带高 = max(行内全部元素字号行高)（防大字号元素压线）
#   · 行数实测：真实字宽折行、每段 ≤ max_lines、符号缩进须已计入 body_w
# 仅几何/算法，不绑配色/字体；各风格传自己的 spacing_range / domains。
# ═══════════════════════════════════════════

def fit_gaps(rows, body_sz, band, concl_sz, container_h, spacing_range, para_n,
             line_spacing=SP_LINE, line_thick=0.014):
    """按给定行数把各间距在区间内智能放大（slack 均摊、cap 各自 max）。

    rows        — 本卡正文总行数
    body_sz     — 正文字号；band — 行带高（max(数字,标题) 行高）
    concl_sz    — 结论字号；container_h — 卡/容器可用高(英寸)
    spacing_range — dict：top/band_to_d1/d1_to_body/para/body_to_d2/d2_to_concl/
                    concl_to_bot，每项 (min, max)（间距下限体现层级）
    para_n      — 段落数-1（段间间距个数）
    返回 gaps dict（各间距实际值，供渲染直接用）。"""
    def lh(sz): return sz * line_spacing / 72.0
    need_min = (spacing_range['top'][0] + band + spacing_range['band_to_d1'][0] + line_thick
                + spacing_range['d1_to_body'][0] + rows * lh(body_sz)
                + para_n * spacing_range['para'][0]
                + spacing_range['body_to_d2'][0] + line_thick
                + spacing_range['d2_to_concl'][0] + lh(concl_sz)
                + spacing_range['concl_to_bot'][0])
    gaps = {k: v[0] for k, v in spacing_range.items()}
    slack = container_h - need_min
    if slack > 0:
        keys = list(spacing_range.keys())
        total_cap = sum(spacing_range[k][1] - spacing_range[k][0] for k in keys)
        if total_cap > 0:
            ratio = min(1.0, slack / total_cap)
            for k in keys:
                room = spacing_range[k][1] - spacing_range[k][0]
                gaps[k] = spacing_range[k][0] + room * ratio
    return gaps

def solve_card_fonts(points_list, body_w, container_h, spacing_range, domains,
                     line_spacing=SP_LINE, max_lines=2):
    """约束求解排版：字号联立求解 + 间距智能弹性（跨风格通用机制）。

    输入：
      points_list   — 各卡段落文本列表的列表，如 [[p1,p2],[p1,p2],...]（并列卡整组同号，以最密卡为约束）
      body_w        — 正文文本宽(英寸)，须已扣除几何符号缩进
      container_h   — 卡/容器可用高(英寸)
      spacing_range — 间距区间 dict（见 fit_gaps）
      domains       — 字号域 dict：body=(min,max) title=(min,max) num=(min,max) concl=(min,max)
      line_spacing  — 段内行距倍数（默认 SP_LINE）
      max_lines     — 每段折行上限（默认 2）
    求解顺序：正文(从上限向下、≥min 底线) → 标题 → 数字 → 结论，取首个「间距取下限时总高 ≤ 容器高」
    的组合；间距按 fit_gaps 弹性放大。无解时返回各角色下限 + 间距全 min。
    返回 (body_sz, title_sz, num_sz, concl_sz, gaps)。"""
    LINE = 0.014
    def lh(sz): return sz * line_spacing / 72.0
    def n_lines(text, sz):
        return max(1, math.ceil(measure_text_width(text, sz) * 1.03 / max(body_w, 0.1)))
    def rows_of(body):
        mx = 0
        for pts in points_list:
            r = 0
            for p in pts:
                n = n_lines(p, body)
                if n > max_lines:
                    return None          # 某段超行数上限 → 该字号不合格
                r += n
            mx = max(mx, r)
        return mx
    para_n = len(points_list[0]) - 1
    for body in range(domains['body'][1], domains['body'][0] - 1, -1):
        mr = rows_of(body)
        if mr is None:
            continue
        for title in range(domains['title'][1], domains['title'][0] - 1, -1):
            if title <= body:
                continue
            for num in range(domains['num'][1], domains['num'][0] - 1, -1):
                if num < title:
                    continue
                band = max(lh(num), lh(title))
                for concl in range(domains['concl'][1], domains['concl'][0] - 1, -1):
                    if concl < body:
                        continue
                    need_min = (spacing_range['top'][0] + band + spacing_range['band_to_d1'][0] + LINE
                                + spacing_range['d1_to_body'][0] + mr * lh(body)
                                + para_n * spacing_range['para'][0]
                                + spacing_range['body_to_d2'][0] + LINE
                                + spacing_range['d2_to_concl'][0] + lh(concl)
                                + spacing_range['concl_to_bot'][0])
                    if need_min > container_h:
                        continue
                    gaps = fit_gaps(mr, body, band, concl, container_h, spacing_range,
                                    para_n, line_spacing)
                    return (body, title, num, concl, gaps)
    gaps = {k: v[0] for k, v in spacing_range.items()}
    return (domains['body'][0], domains['title'][0], domains['num'][0],
            domains['concl'][0], gaps)

# ═══════════════════════════════════════════
# 布局骨架 × 皮肤（架构层，2026-07-31 用户拍板：布局统一向企业风对齐、四风格自动生效）
# 原则：布局几何只写一份（本层）——企业风调通 = 骨架调通；四风格只提供 skin（色/字体/字号/箭头/卡样式）。
# skin 契约字段（各引擎实现 SKIN_* dict）：
#   bg          — 页面底色
#   fonts       — {head: 标题字体, body: 正文字体}
#   cap         — {size, bold, color} 编号小标题（①②③ 引导层）
#   node        — {size, bold, color, sub_size, sub_color} 闭环节点
#   card        — {title_size, title_color, sub_size, sub_color, body_size, body_color,
#                  concl_color, line, line_w} 双卡（含 D1 线/结论尾句色）
#   arrow       — {color, w} 箭头（色随底色、浅底细线）
#   tail_bar    — 第三块说明文本的悬挂引注色（None = 不画）
#   line_spacing— 行距倍数（默认 1.2）
# ═══════════════════════════════════════════

def body_hl(tf, text, hl, size=12, color=LIGHT_GRAY, hl_color=BIZ_TEAL, fonts=None, sep="。"):
    """正文 + 加粗强调尾句（同段混合 run，供布局骨架结论尾句用）。"""
    body_f = fonts.get('body', FONT) if fonts else FONT
    p = tf.paragraphs[0]
    for t, bold, c in ((text, False, color), (sep + hl, True, hl_color)):
        if not t:
            continue
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = body_f
        r.font.color.rgb = c
        set_chinese_font(r, body_f)

def _skin_text(slide, x, y, w, h, text, size, bold, color, font, align=PP_ALIGN.LEFT):
    """按皮肤字体渲染一段文本（骨架内部用）。"""
    tf = tb_box(slide, x, y, w, h)
    p = run_text(tf, text, size, bold, color, align, Pt(0), True)
    for r in p.runs:
        r.font.name = font
        set_chinese_font(r, font)
    return tf


def _bullet_text(slide, x, y, w, h, text, kw, size, color, kw_color, font):
    """要点文本：可选核心词高亮（kw 为 str 或 (word,...) 元组，命中段 bold+强调色，克制每段 1~2 词）。"""
    tf = tb_box(slide, x, y, w, h)
    if not kw:
        p = run_text(tf, text, size, False, color, PP_ALIGN.LEFT, Pt(0), True)
        for r in p.runs:
            r.font.name = font; set_chinese_font(r, font)
        return tf
    kws = [kw] if isinstance(kw, str) else list(kw)
    segs = [text]
    for w0 in kws:
        if not w0:
            continue
        nxt = []
        for s0 in segs:
            if w0 in s0:
                head, tail = s0.split(w0, 1)
                if head:
                    nxt.append(head)
                nxt.append((w0, True))
                if tail:
                    nxt.append(tail)
            else:
                nxt.append(s0)
        segs = nxt
    p = tf.paragraphs[0]
    for seg in segs:
        if isinstance(seg, tuple):
            r = p.add_run(); r.text = seg[0]
            r.font.size = Pt(size); r.font.bold = True; r.font.name = font
            r.font.color.rgb = kw_color; set_chinese_font(r, font)
        else:
            r = p.add_run(); r.text = seg
            r.font.size = Pt(size); r.font.name = font
            r.font.color.rgb = color; set_chinese_font(r, font)
    return tf

def layout_loop_page(slide, skin, content, y0=1.88, y_end=6.45,
                     x0=CANVAS_X0, cw=CANVAS_CW, card_h=None, tail_h=None,
                     stretch=True):
    """闭环复合页公共骨架（跨风格统一布局，企业风样板节奏）。

    结构：① 编号cap + 闭环节点（name/sub、顶行箭头、U型回流）
         → ② 编号cap + 双卡（title/sub → D1 → desc+结论尾句）+ ⇄ 协同
         → ③ 编号cap + 第三块（tail：卡列表 或 说明文本列表）
    content: {cap1, nodes, cap2, cards, cap3, tail, contrast?, bg_caption?}
      nodes = [{name, sub?}, ...×4]；cards = [{title, sub?, desc, concl?}, ...×2]
      tail = [{title, desc}, ...]（卡）或 [str, ...]（说明文本）
    card_h / tail_h：②双卡 / ③卡列表 高度（None = 自动 1.44 / 0.86，文质节奏）；
      传值且 cards 无 sub/concl 时走紧凑两行（title+desc，企业 work_people 角色/能力卡节奏）。
    页眉 section_header 与底部金句由调用方处理（本函数只管内容区三块）。
    返回实际底部 y。"""
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    cap, nd, cd, ar = skin['cap'], skin['node'], skin['card'], skin['arrow']
    head_f = skin['fonts']['head']
    n0 = len(slide.shapes)                 # 记录调用前形状数，便于整体垂直居中
    # ── 间距节奏（固定基线 + 可选拉伸）──
    cap_gap = lh(cap['size']) + 0.12          # cap → 内容（小间距，固定，保层级）
    loop_pad = 0.22                            # 闭环底部 U 回流留白（固定）
    base_block_gap = 0.14                      # 块 → 下个 cap（可拉伸主间隙）
    nw, ng = 2.55, (cw - 4 * 2.55) / 3.0
    has_sub = any(n.get('sub') for n in content['nodes'])
    nh = 0.72 if has_sub else 0.55
    card_h = card_h if card_h is not None else 1.44
    tail = content.get('tail') or []
    th = tail_h if tail_h is not None else 0.86
    contrast_extra = 0.34 if content.get('contrast') else 0.0
    bg_extra = 0.32 if content.get('bg_caption') else 0.0
    # 固定总高（不含拉伸），用于把余量均摊到块间间隙
    H_fixed = (3 * cap_gap
               + (nh + loop_pad + base_block_gap)
               + (card_h + base_block_gap)
               + (th + base_block_gap)
               + contrast_extra + bg_extra)
    block_gap = base_block_gap
    if stretch and y_end > y0:
        surplus = (y_end - y0) - H_fixed
        if surplus > 0.02:
            # 余量均摊到 3 个块间间隙；单间隙上限 0.45，超出部分仍交 _vshift 上下居中
            per = min(surplus / 3.0, 0.45)
            block_gap = base_block_gap + per
    # ── ① 编号cap + 闭环 ──
    y = y0
    _skin_text(slide, x0, y, cw, 0.26, content['cap1'], cap['size'], cap.get('bold', True),
               cap['color'], head_f)
    y += cap_gap
    centers = []
    for i, n in enumerate(content['nodes']):
        nx = x0 + i * (nw + ng)
        content_card(slide, nx, y, nw, nh, cd['line'], line_w=cd.get('line_w', 1.25))
        if has_sub:
            _skin_text(slide, nx + 0.12, y + 0.08, nw - 0.24, 0.28, n['name'],
                       nd['size'], nd.get('bold', True), nd['color'], head_f, PP_ALIGN.CENTER)
            _skin_text(slide, nx + 0.12, y + 0.40, nw - 0.24, 0.24, n.get('sub', ''),
                       nd['sub_size'], False, nd['sub_color'], skin['fonts']['body'], PP_ALIGN.CENTER)
        else:
            tf = tb_box(slide, nx + 0.12, y, nw - 0.24, nh)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = run_text(tf, n['name'], nd['size'], nd.get('bold', True), nd['color'],
                         PP_ALIGN.CENTER, Pt(0), True)
            for r in p.runs:
                r.font.name = head_f; set_chinese_font(r, head_f)
        centers.append(nx + nw / 2.0)
        if i < 3:
            loop_arrow(slide, nx + nw + 0.09, y + nh / 2.0, nx + nw + ng - 0.09, y + nh / 2.0,
                       color=ar['color'], w=ar['w'])
    loop_return_u(slide, centers, y + nh, y + nh + 0.22, color=ar['color'], w=ar['w'])
    y += nh + loop_pad + block_gap
    if content.get('contrast'):                       # ① 区后小字（企业风 work_people 可选）
        _skin_text(slide, x0, y, cw, 0.26, content['contrast'], 12, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f)
        y += 0.34
    # ── ② 编号cap + 双卡 + ⇄ ──
    _skin_text(slide, x0, y, cw, 0.26, content['cap2'], cap['size'], cap.get('bold', True),
               cap['color'], head_f)
    y += cap_gap
    cw2 = (cw - 0.33) / 2.0
    body_w = cw2 - 0.56
    compact = card_h < 1.1            # 紧凑两行模式（企业 work_people 角色卡：title+desc）
    for i, c in enumerate(content['cards']):
        x = x0 + i * (cw2 + 0.33)
        content_card(slide, x, y, cw2, card_h, cd['line'], line_w=cd.get('line_w', 1.25))
        _skin_text(slide, x + 0.28, y + 0.14, body_w, 0.30, c['title'],
                   cd['title_size'], True, cd['title_color'], head_f)
        if c.get('sub'):
            _skin_text(slide, x + 0.28, y + 0.47, body_w, 0.22, c['sub'],
                       cd['sub_size'], False, cd['sub_color'], skin['fonts']['body'])
        if compact:
            _skin_text(slide, x + 0.28, y + 0.52, body_w, 0.30, c['desc'],
                       cd['body_size'], False, cd['body_color'], skin['fonts']['body'])
            continue
        d1_y = y + 0.79
        card_divider(slide, x + 0.28, d1_y, body_w, cd['line'], pt=0.75, full=False,
                     title_w=min(measure_text_width(c['title'], cd['title_size']) + 0.24, body_w))
        tf = tb_box(slide, x + 0.28, d1_y + 0.13, body_w, 0.44)
        if c.get('concl'):
            body_hl(tf, c['desc'], c['concl'], size=cd['body_size'],
                    color=cd['body_color'], hl_color=cd['concl_color'], fonts=skin['fonts'])
        else:
            p = run_text(tf, c['desc'], cd['body_size'], False, cd['body_color'],
                         PP_ALIGN.LEFT, Pt(0), True)
            for r in p.runs:
                r.font.name = skin['fonts']['body']; set_chinese_font(r, skin['fonts']['body'])
    tf = tb_box(slide, x0 + cw2 + 0.03, y + card_h / 2 - 0.25, 0.34, 0.50)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = run_text(tf, "⇄", 22, True, cd['title_color'], PP_ALIGN.CENTER, Pt(0), True)
    for r in p.runs:
        r.font.name = skin['fonts']['body']; set_chinese_font(r, skin['fonts']['body'])
    y += card_h + block_gap
    if content.get('bg_caption'):                     # ② 区后小字（企业风 work_people 可选）
        _skin_text(slide, x0, y, cw, 0.24, content['bg_caption'], 12, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f)
        y += 0.32
    # ── ③ 编号cap + 第三块 ──
    _skin_text(slide, x0, y, cw, 0.26, content['cap3'], cap['size'], cap.get('bold', True),
               cap['color'], head_f)
    y += cap_gap
    if tail and isinstance(tail[0], dict):            # 卡列表（结构化）
        tw = (cw - 2 * 0.30) / 3.0
        compact_tail = th < 0.7                       # 紧凑能力卡（企业 work_people 第三块）
        for i, t in enumerate(tail):
            x = x0 + i * (tw + 0.30)
            content_card(slide, x, y, tw, th, cd['line'], line_w=cd.get('line_w', 1.25))
            if compact_tail:
                _skin_text(slide, x + 0.20, y + 0.06, tw - 0.40, 0.26, t['title'],
                           14, True, cd['title_color'], head_f)
                _skin_text(slide, x + 0.20, y + 0.32, tw - 0.40, 0.26, t['desc'],
                           12, False, cd['body_color'], skin['fonts']['body'])
            else:
                _skin_text(slide, x + 0.20, y + 0.10, tw - 0.40, 0.26, t['title'],
                           cd['title_size'], True, cd['title_color'], head_f)
                _skin_text(slide, x + 0.20, y + 0.40, tw - 0.40, 0.40, t['desc'],
                           cd['body_size'], False, cd['body_color'], skin['fonts']['body'])
        y += th + block_gap
    else:                                             # 说明文本列表（悬挂引注）
        bar = skin.get('tail_bar')
        for t in tail:
            tx, tw2 = (x0 + 0.16, cw - 0.16) if bar else (x0, cw)
            h = max(1, math.ceil(measure_text_width(t, 12) * 1.03 / tw2)) * 12 * ls / 72.0
            if bar:
                rect(slide, x0, y + 0.02, 0.022, h, fill=bar, line=None)
            tf = tb_box(slide, tx, y, tw2, h + 0.08)
            p = run_text(tf, t, 12, False, skin['fonts'].get('mute', cd['body_color']),
                         PP_ALIGN.LEFT, Pt(0), True)
            for r in p.runs:
                r.font.name = skin['fonts']['body']; set_chinese_font(r, skin['fonts']['body'])
            y += h + 0.02
    _vshift(slide, n0, y0, y, y_end)       # 内容块在内容区垂直居中（消顶重）
    return y

def layout_timeline(slide, skin, items, y0, y_end, x0=CANVAS_X0, cw=CANVAS_CW):
    """时间轴公共骨架：竖线 + 圆点 + 日期(右对齐) + 事件文本，整体垂直居中。

    items = [{date, text}, ...]（日期走 fmt_date_5digit 统一 5 位）；
    企业风样板：竖线底部 tailEnd 箭头可选（皮肤 timeline.arrow）。
    皮肤 timeline 字段：axis(轴色) / dot(圆点色) / date_size / date_color / event_size /
    event_color / arrow(布尔，企业风 True)。"""
    tl = skin['timeline']
    axis_x = x0 + 0.80                       # 轴 x = 1.60（企业/文质同值）
    top, bottom = y0, y_end
    rect(slide, axis_x - 0.01, top, 0.02, bottom - top, fill=tl['axis'], line=None)
    if tl.get('arrow'):                       # 企业风：竖线底部 tailEnd 箭头（线宽走箭头约定 ≈2pt）
        cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            Inches(axis_x - 0.01), Inches(bottom - 0.14), Inches(axis_x - 0.01), Inches(bottom - 0.02))
        cxn.line.color.rgb = tl['axis']
        cxn.line.width = Pt(tl.get('arrow_w', 2))
        ln = cxn.line._get_or_add_ln()
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    step = (bottom - top) / max(1, len(items))
    body_f = skin['fonts']['body']
    for i, it in enumerate(items):
        y = top + step * i + step / 2
        rect(slide, axis_x - 0.07, y - 0.07, 0.14, 0.14, fill=tl['dot'], line=None, rounded=True)
        _dw = axis_x - x0 - 0.25
        _dh = measure_text_height(fmt_date_5digit(it['date']), tl['date_size'], _dw, 1.2) + 0.04
        tf = tb_box(slide, x0, y - _dh / 2.0, _dw, _dh)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, fmt_date_5digit(it['date']), tl['date_size'], True, tl['date_color'],
                     PP_ALIGN.RIGHT, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
        _ew = cw - (axis_x - x0) - 0.40
        # 金额（amt）：追加琥珀加粗后缀 " ｜ 40亿$"，落实"金额高亮"DNA（企业风走专属徽章，此处覆盖其余三风）
        disp = it['text']
        kw_a = list(it.get('kw_amber') or [])
        if it.get('amt'):
            disp = it['text'] + "　｜　" + it['amt']
            kw_a.append(it['amt'])
        kw_t = it.get('kw_teal') or []
        _eh = measure_text_height(disp, tl['event_size'], _ew, 1.2) + 0.04
        tf = tb_box(slide, axis_x + 0.40, y - _eh / 2.0, _ew, _eh)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        emph_runs(tf, disp, tl['event_size'], tl['event_color'],
                   kw_t, kw_a,
                   color_t=tl.get('kw_teal_color', tl['dot']),
                   color_a=tl.get('kw_amber_color', tl['event_color']),
                   font=body_f)
    return bottom

def layout_card_grid(slide, skin, cards, cols=2, y0=1.88, y_end=6.45,
                     x0=CANVAS_X0, cw=CANVAS_CW):
    """卡片网格公共骨架（2×2 通用）：等宽卡，**框高=内容自适应(紧抱内容)+整组垂直居中**，消死白。

    设计取向（对齐 ChatGPT 参考 + 用户拍板）：**小字 + 大留白 + 颜色拉层级**，不顶满卡片。
    原"取最大能装字号"求解器已删除（用户反馈"字体太大/紧凑"）；后又删 auto-fit 13（13/14 浮动是隐形不一致）。
    卡内节奏（对齐 SKILL §4.3.3「五区」）：
      ① 序号 tag（固定 22pt 粗·强调色·左上，全场最大作主标识）与 标题 title（固定 14pt 粗·标题色）
         **同行、MIDDLE 垂直居中**（序号占 NUM_W 独立框、已收窄避免推远标题，标题紧随其后）；
      ② 正文 lines（固定 12pt 灰·左对齐·"·" 前缀 bullet·行距 1.35 宽松·段后 0.10″ 呼吸）；
      ③ 分隔线（隐约暗色退背景层，正文↔结论间，跟在正文之后，不写死到卡底）；
      ④ 结论 hl（固定 14pt 强调色·居中·加粗，**锁死 14、禁 auto-fit 13**；受 _guard_one_line 硬规则约束**禁止换行**）。
    框高机制：ch = band_h(c) + 2*VPAD（内容实测高 + 上下各 VPAD 内边距 → 框紧抱内容，消死白）；
      整组 2×N 在 [y0,y_end] 内垂直居中，框间距 gap 固定为统一呼吸 → 内部节奏与框间距同频、视觉协调。
    红线（§4.12-7）：序号须 22pt（非正文 12pt）；结论须居中+加粗+**单行禁换行**（14pt 锁死）；
      正文带 "·" 前缀；整组距金句线 ≥0.10″、卡底 < 金句线 6.50。
    字号现为布局标准化固定值（覆盖皮肤 card 的 title_size 等），颜色仍取自皮肤。"""
    cd = skin['card']
    body_f = skin['fonts']['body']
    head_f = skin['fonts']['head']
    num_color = cd.get('num_color', cd.get('concl_color', cd['title_color']))  # 编号可独立于结论取色（商务风：编号金/结论深蓝）
    ls = skin.get('line_spacing', 1.2)
    # ── 固定舒适小字号 + 大留白（对齐 ChatGPT 参考：小字、颜色拉层级、不顶满）──
    # 删除原"取最大能装字号"求解器：那会顶满卡片、挤压上下留白（用户反馈"字体太大/紧凑"）。
    NUM_SZ = 22          # 序号：全场最大+强调色+粗
    TS = 14              # 标题 14pt
    BS = 12              # 正文 12pt（SKILL 底线，不可再小）
    CONCL_SZ = 14        # 结论 14pt（锁死，禁 auto-fit 13——13/14 浮动本身是隐形不一致）
    BODY_LS = 1.35       # 正文行距固定宽松
    BULLET_GAP = 0.10    # bullet 段后呼吸
    pad = 0.24           # 卡内左右内边距
    VPAD = 0.16          # 卡内上下内边距（框紧抱内容，消死白）
    rows = (len(cards) + cols - 1) // cols
    gap = 0.32           # 框间距（统一呼吸，与内部 VPAD 同频）
    cw2 = (cw - gap * (cols - 1)) / cols
    iw = cw2 - 2 * pad
    NUM_W = 0.52         # 序号占位宽（原 0.78 过宽→推远标题+挤窄标题，收窄）
    NUM_TITLE_GAP = 0.12 # 序号→标题间距（原 0.16）
    BAND_GAP = 0.10      # 序号/标题行→正文间距
    DIV_GAP = 0.10       # 正文→分隔线
    CONCL_GAP = 0.10     # 分隔线→结论

    def bullet_rows(ln):
        return max(1, math.ceil(measure_text_width("· " + ln, BS) * 1.04 / iw))

    def band_h(c):
        """内容带真实高度（固定字号下）；框高 = band_h + 2*VPAD → 框紧抱内容，消死白。"""
        bh = 0.0
        for ln in c.get('lines', []):
            bh += bullet_rows(ln) * BS * BODY_LS / 72.0 + BULLET_GAP
        if c.get('lines'):
            bh -= BULLET_GAP
        concl_h = CONCL_SZ * ls / 72.0 + 0.04   # 结论锁 14pt（禁 auto-fit 13）
        return (NUM_SZ * ls / 72.0 + 0.04) + BAND_GAP + bh + DIV_GAP + 0.01 + CONCL_GAP + concl_h

    # ── 框高 = 内容自适应；整组垂直居中于 [y0,y_end]（框紧抱内容 + 间距同频，视觉协调）──
    avail = (y_end - y0)
    ch = band_h(cards[0]) + 2 * VPAD          # 各卡结构相同 → 框高一致
    block_h = rows * ch + (rows - 1) * gap
    sy = y0 + max(0.0, (avail - block_h) / 2.0)
    for i, c in enumerate(cards):
        r, cc = divmod(i, cols)
        x = x0 + cc * (cw2 + gap)
        y = sy + r * (ch + gap)
        content_card(slide, x, y, cw2, ch, cd['line'], line_w=cd.get('line_w', 1.25),
                     shadow_alpha=cd.get('shadow_alpha', 8000),
                     shadow_blur=cd.get('shadow_blur', 90000),
                     shadow_dist=cd.get('shadow_dist', 20000))
        ix = x + pad
        top = y + VPAD                             # 内容自框内上边距起排（紧抱，非居中留死白）
        band = NUM_SZ * ls / 72.0 + 0.04
        # ── 金句硬规则：卡内结论强制单行、禁换行（SKILL §4.12-7）──
        _guard_one_line(c['hl'], CONCL_SZ, iw, ctx=f"卡{c.get('tag','?')}结论")
        # ① 序号 + 标题 同行、MIDDLE 垂直居中
        tf = tb_box(slide, ix, top, NUM_W, band)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, c.get('tag', ''), NUM_SZ, True,
                     num_color, PP_ALIGN.LEFT, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = head_f; set_chinese_font(r2, head_f)
        tf = tb_box(slide, ix + NUM_W + NUM_TITLE_GAP, top,
                    iw - NUM_W - NUM_TITLE_GAP, band)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, c['title'], TS, True, cd['title_color'], PP_ALIGN.LEFT, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = body_f; set_chinese_font(r2, body_f)
        # ② 正文 bullet（"· " 前缀 + 宽松行距 + 呼吸段距）
        yy = top + band + BAND_GAP
        for ln in c.get('lines', []):
            nr = bullet_rows(ln)
            lh = nr * BS * BODY_LS / 72.0
            tf = tb_box(slide, ix, yy, iw, lh + 0.04)
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = run_text(tf, "· " + ln, BS, False, cd['body_color'],
                         PP_ALIGN.LEFT, Pt(0), True)
            for r2 in p.runs:
                r2.font.name = body_f; set_chinese_font(r2, body_f)
            yy += lh + BULLET_GAP
        # ③ 分隔线（跟在正文之后，不写死到卡底）
        div_y = yy + DIV_GAP
        card_divider(slide, ix, div_y, iw, cd['line'], pt=0.75, full=True)
        # ④ 结论（居中 + 加粗 + 锁 14pt，禁 auto-fit 13；受 _guard_one_line 约束禁换行）
        concl_h = CONCL_SZ * ls / 72.0 + 0.04
        tf = tb_box(slide, ix, div_y + CONCL_GAP, iw, concl_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, c['hl'], CONCL_SZ, True, cd['concl_color'],
                     PP_ALIGN.CENTER, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = body_f; set_chinese_font(r2, body_f)
    return y_end


def layout_compare(slide, skin, content, y0=1.85, y_end=6.30,
                   x0=CANVAS_X0, cw=CANVAS_CW):
    """定义对比公共骨架（企业风 definition_compare 样板）：左宽定义区 + 右窄排除区 + 中部竖线。

    content = {hero, hero_en, hero_cn, bullets[], summary?, right:{header, negatives[], positive}}
      negatives = [{kw, desc}]；positive = {kw, desc}（✓ 高亮框）
    skin['compare'] 字段：
      hero_size/hero_color · en_size/en_color/en_italic · cn_size/cn_color
      bullet_size/bullet_color/bullet_mark("dot"|"bar")/mark_color
      summary_size/summary_color
      header_size/header_color
      neg_badge("circle"|"text")/neg_kw_size/neg_kw_color/neg_desc_size/neg_desc_color
        /neg_badge_bg/neg_badge_color
      pos_bg/pos_border/pos_kw_size/pos_kw_color/pos_desc_size/pos_desc_color
      div(竖线色)/right_card(右区画卡框)/inner_div(卡内分隔线色)
    页眉与底部金句由调用方处理。"""
    cp = skin.get('compare', {})
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    RIGHT_CARD = cp.get('right_card', False)
    DIV_X = x0 + 7.20                    # 企业风竖线 x=7.80（x0=0.60）
    left_w = (DIV_X - x0 - 0.30) if not RIGHT_CARD else 6.30
    if RIGHT_CARD:
        right_x = x0 + left_w + 0.35
    else:
        right_x = DIV_X + 0.20
    right_w = cw + x0 - right_x - 0.20
    # ── 中部竖线（无卡框模式）／卡框（右区）──
    if not RIGHT_CARD:
        rect(slide, DIV_X, y0, Pt(1.25), y_end - y0, fill=cp.get('div', skin['card']['line']))
    else:
        content_card(slide, right_x, y0, right_w, y_end - y0,
                     skin['card']['line'], line_w=skin['card'].get('line_w', 1.25), adj=0.03)
    # ── 左区：hero + en + cn + bullets + summary ──
    y = y0 + 0.10
    _skin_text(slide, x0, y, left_w, lh(cp['hero_size']) + 0.05, content['hero'],
               cp['hero_size'], True, cp['hero_color'], head_f)
    y += lh(cp['hero_size']) + 0.02
    tf = tb_box(slide, x0, y, left_w, lh(cp['en_size']) + 0.05)
    p = run_text(tf, content['hero_en'], cp['en_size'], False, cp['en_color'],
                 PP_ALIGN.LEFT, Pt(0), True)
    if cp.get('en_italic', True):
        p.runs[0].font.italic = True
    for r in p.runs:
        r.font.name = body_f; set_chinese_font(r, body_f)
    y += lh(cp['en_size']) + 0.04
    _skin_text(slide, x0, y, left_w, lh(cp['cn_size']) + 0.05, content['hero_cn'],
               cp['cn_size'], True, cp['cn_color'], head_f)
    y += lh(cp['cn_size']) + 0.32
    mark = cp.get('bullet_mark', 'dot')
    kw_color = cp.get('bullet_kw_color', cp.get('mark_color', cp['bullet_color']))
    for b in content['bullets']:
        text, kw = (b[0], b[1]) if isinstance(b, (tuple, list)) else (b, None)
        h = max(1, math.ceil(measure_text_width(text, cp['bullet_size']) * 1.03 / left_w)) \
            * lh(cp['bullet_size'])
        if mark == 'dot':
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x0 + 0.02),
                Inches(y + h / 2 - 0.05), Inches(0.10), Inches(0.10))
            dot.fill.solid(); dot.fill.fore_color.rgb = cp.get('mark_color', cp['bullet_color'])
            dot.line.fill.background(); dot.shadow.inherit = False
            _bullet_text(slide, x0 + 0.25, y, left_w - 0.25, h + 0.04, text, kw,
                         cp['bullet_size'], cp['bullet_color'], kw_color, body_f)
        else:                                   # bar 悬挂引注（文质风）
            rect(slide, x0, y + 0.02, 0.022, h, fill=cp.get('mark_color', cp['bullet_color']))
            _bullet_text(slide, x0 + 0.16, y, left_w - 0.16, h + 0.04, text, kw,
                         cp['bullet_size'], cp['bullet_color'], kw_color, body_f)
        y += h + 0.16
    if content.get('summary'):
        y += 0.10
        _skin_text(slide, x0, y, left_w, lh(cp['summary_size']) + 0.05,
                   content['summary'], cp['summary_size'], True, cp['summary_color'], head_f)
    # ── 右区：header + negatives + positive ──
    rx = right_x + (0.30 if RIGHT_CARD else 0.0)
    rw = right_w - (0.60 if RIGHT_CARD else 0.0)
    badge = cp.get('neg_badge', 'circle')
    # 右栏内容垂直居中于卡框内（RIGHT_CARD 模式）：避免卡底留大片空白，违反 §4.4
    if RIGHT_CARD:
        card_h = y_end - y0
        _rh = lh(cp['header_size']) + 0.30
        for it in content['right']['negatives']:
            if badge == 'circle':
                kw_w = measure_text_width(it['kw'], cp['neg_kw_size'])
                desc_w = measure_text_width(it['desc'], cp['neg_desc_size']) + 0.15
                rows_n = max(1, math.ceil((kw_w + desc_w) * 1.03 / (rw - 0.50)))
                h = rows_n * lh(cp['neg_kw_size']) + 0.05
            else:
                total_w = measure_text_width("✕　%s —— " % it['kw'], cp['neg_kw_size']) \
                          + measure_text_width(it['desc'], cp['neg_desc_size'])
                rows_n = max(1, math.ceil(total_w * 1.03 / rw))
                h = rows_n * lh(cp['neg_kw_size']) + 0.05
            _rh += h + (0.18 if badge == 'circle' else 0.16)
        if badge == 'circle':
            pos_h = lh(cp['pos_kw_size']) + lh(cp['pos_desc_size']) + 0.34
        else:
            desc_text = content['right']['positive']['desc']
            rows_n = max(1, math.ceil(measure_text_width(desc_text, cp['pos_desc_size']) * 1.03 / rw))
            desc_h = rows_n * lh(cp['pos_desc_size']) + 0.05
            pos_h = lh(cp['pos_kw_size']) + 0.02 + desc_h
        _rh += 0.22 + pos_h
        y = y0 + max(0.12, (card_h - _rh) / 2.0)
    else:
        y = y0 + 0.10
    _skin_text(slide, rx, y, rw, lh(cp['header_size']) + 0.05, content['right']['header'],
               cp['header_size'], True, cp['header_color'], head_f)
    y += lh(cp['header_size']) + 0.30
    for it in content['right']['negatives']:
        if badge == 'circle':
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(rx + 0.02),
                Inches(y + 0.02), Inches(0.34), Inches(0.34))
            circ.fill.solid(); circ.fill.fore_color.rgb = cp['neg_badge_bg']
            circ.line.fill.background(); circ.shadow.inherit = False
            ctf = circ.text_frame
            for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
                setattr(ctf, m, 0)
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            cpar = ctf.paragraphs[0]; cpar.alignment = PP_ALIGN.CENTER
            cr = cpar.add_run(); cr.text = "×"
            cr.font.size = Pt(16); cr.font.bold = True; cr.font.color.rgb = cp['neg_badge_color']
            cr.font.name = body_f; set_chinese_font(cr, body_f)
            kw_w = measure_text_width(it['kw'], cp['neg_kw_size'])
            desc_w = measure_text_width(it['desc'], cp['neg_desc_size']) + 0.15
            rows_n = max(1, math.ceil((kw_w + desc_w) * 1.03 / (rw - 0.50)))
            h = rows_n * lh(cp['neg_kw_size']) + 0.05
            tf = tb_box(slide, rx + 0.50, y, rw - 0.50, h)
            tf.word_wrap = True
            pp = tf.paragraphs[0]
            kr = pp.add_run(); kr.text = it['kw']
            kr.font.size = Pt(cp['neg_kw_size']); kr.font.bold = True
            kr.font.color.rgb = cp['neg_kw_color']; kr.font.name = body_f; set_chinese_font(kr, body_f)
            dr = pp.add_run(); dr.text = "  " + it['desc']
            dr.font.size = Pt(cp['neg_desc_size']); dr.font.color.rgb = cp['neg_desc_color']
            dr.font.name = body_f; set_chinese_font(dr, body_f)
            y += h + 0.18
        else:                                   # text 前缀（文质风）
            total_w = measure_text_width("✕　%s —— " % it['kw'], cp['neg_kw_size']) \
                      + measure_text_width(it['desc'], cp['neg_desc_size'])
            rows_n = max(1, math.ceil(total_w * 1.03 / rw))
            h = rows_n * lh(cp['neg_kw_size']) + 0.05
            tf = tb_box(slide, rx, y, rw, h)
            tf.word_wrap = True
            pp = tf.paragraphs[0]
            kr = pp.add_run(); kr.text = "✕　%s —— " % it['kw']
            kr.font.size = Pt(cp['neg_kw_size']); kr.font.bold = True
            kr.font.color.rgb = cp['neg_kw_color']; kr.font.name = body_f; set_chinese_font(kr, body_f)
            dr = pp.add_run(); dr.text = it['desc']
            dr.font.size = Pt(cp['neg_desc_size']); dr.font.color.rgb = cp['neg_desc_color']
            dr.font.name = body_f; set_chinese_font(dr, body_f)
            y += h + 0.16
        if it is content['right']['negatives'][-1] and cp.get('inner_div'):
            d_y = y - 0.09
            rect(slide, rx, d_y, rw, 0.02, fill=cp['inner_div'])
    # ✓ 正解高亮框（企业风青底圆角 / 文质风透明青边）
    y += 0.22
    pos_h = lh(cp['pos_kw_size']) + lh(cp['pos_desc_size']) + 0.34
    if badge == 'circle':
        box = rounded_card(slide, right_x, y, right_w, pos_h,
                           fill=cp.get('pos_bg'), adj=0.05,
                           line=cp.get('pos_border'), line_w=1.0)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(right_x + 0.15),
            Inches(y + pos_h / 2 - 0.17), Inches(0.34), Inches(0.34))
        circ.fill.solid(); circ.fill.fore_color.rgb = cp['pos_border']
        circ.line.fill.background(); circ.shadow.inherit = False
        ctf = circ.text_frame
        for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
            setattr(ctf, m, 0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cpar = ctf.paragraphs[0]; cpar.alignment = PP_ALIGN.CENTER
        cr = cpar.add_run(); cr.text = "✓"
        cr.font.size = Pt(14); cr.font.bold = True; cr.font.color.rgb = cp.get('neg_badge_bg', WHITE)
        cr.font.name = body_f; set_chinese_font(cr, body_f)
        tf = tb_box(slide, right_x + 0.60, y + 0.08, right_w - 0.75, pos_h - 0.16)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]
        if content['right']['positive'].get('kw'):
            kr = pp.add_run(); kr.text = content['right']['positive']['kw']
            kr.font.size = Pt(cp['pos_kw_size']); kr.font.bold = True
            kr.font.color.rgb = cp['pos_kw_color']; kr.font.name = body_f; set_chinese_font(kr, body_f)
            kr = pp.add_run(); kr.text = "  "
            kr.font.size = Pt(cp['pos_kw_size']); kr.font.color.rgb = cp['pos_kw_color']
            kr.font.name = body_f; set_chinese_font(kr, body_f)
        dr = pp.add_run(); dr.text = content['right']['positive']['desc']
        dr.font.size = Pt(cp['pos_desc_size']); dr.font.color.rgb = cp['pos_desc_color']
        dr.font.name = body_f; set_chinese_font(dr, body_f)
    else:
        _skin_text(slide, rx, y, rw, lh(cp['pos_kw_size']) + 0.05,
                   "✓　" + content['right']['positive']['kw'],
                   cp['pos_kw_size'], True, cp['pos_kw_color'], head_f)
        y += lh(cp['pos_kw_size']) + 0.02
        # 正解描述：按内容行数自适应高度 + 自动折行（避免单线框被撑破换行）
        desc_text = content['right']['positive']['desc']
        rows_n = max(1, math.ceil(measure_text_width(desc_text, cp['pos_desc_size']) * 1.03 / rw))
        desc_h = rows_n * lh(cp['pos_desc_size']) + 0.05
        tf = tb_box(slide, rx, y, rw, desc_h)
        tf.word_wrap = True
        p = run_text(tf, desc_text, cp['pos_desc_size'], False, cp['pos_desc_color'],
                     PP_ALIGN.LEFT, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
    return y_end


def layout_value_grid(slide, skin, items, y0=1.80, y_end=6.48,
                      x0=CANVAS_X0, cw=CANVAS_CW, loop=True):
    """2×2 价值网格公共骨架（企业风 value_loop 样板：编号+标题同行 → 要点 → 分隔线 → 结论）。

    items = [{num, title, points[], tag}, ×4]（左上→右上→右下→左下）
    loop=True → 四段等长循环箭头（企业风：四卡上下左右间距全等 gap=gap_v=gap_h，skin 覆盖为 0.5，白/深箭头）
    loop=False → 纯并列（文质风：gap_h 0.33 / gap_v 0.16，无箭头）
    skin['value_grid'] 字段：
      num_size/num_color · title_size/title_color · body_size/body_color
      tag_size/tag_color/tag_align("center"|"left")
      div(分隔线色)/d1(布尔：标题下 D1 短线)/mark(布尔：段落青方块)
      gap_h/gap_v · solve(布尔：True 走公共求解器 ≥12 底线+间距弹性 / False 企业风固定节奏)
      line/line_w（卡框，缺省回退 skin['card']）
    卡内六区统一：带(num+title 同行，行带=max 行高) → [D1] → 要点 → D2(分隔线) → 结论。"""
    cd = skin['card']
    vg = skin.get('value_grid', {})
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    line = vg.get('line', cd['line']); line_w = vg.get('line_w', cd.get('line_w', 1.25))
    gap_h = vg.get('gap_h', 0.70 if loop else 0.33)
    gap_v = vg.get('gap_v', 0.70 if loop else 0.16)
    # 卡宽策略：measured（企业风定稿：正文最长行实测+buffer，整簇水平居中）/ even（均分）
    if vg.get('card_w') == 'measured':
        longest = 0.0
        for it in items:
            for ln in it['points']:
                longest = max(longest, measure_text_width("· " + ln, 12))
        cw2 = max(3.0, round((longest + 0.60 + 0.15) / 0.05) * 0.05)
        if vg.get('center', False):
            x0 = (13.333 - (2 * cw2 + gap_h)) / 2   # 左右留白相等（企业风定稿）
    else:
        cw2 = (cw - gap_h) / 2.0
    ch = (y_end - y0 - gap_v) / 2.0
    PAD = 0.28
    body_w = cw2 - 2 * PAD
    # ── 字号：solve=True 走公共求解器（正文 ≥12 底线、间距区间弹性）──
    if vg.get('solve', False):
        SP_RANGE = vg.get('sp_range', dict(top=(0.08, 0.12), band_to_d1=(0.12, 0.16),
                         d1_to_body=(0.12, 0.16), para=(0.10, 0.14), body_to_d2=(0.12, 0.16),
                         d2_to_concl=(0.10, 0.14), concl_to_bot=(0.06, 0.10)))
        DOMAINS = vg.get('domains', dict(body=(12, 16), title=(14, 18),
                        num=(16, 20), concl=(11, 14)))
        MARK_W = 0.16 if vg.get('mark', False) else 0.0
        body_w = body_w - MARK_W
        body_sz, title_sz, num_sz, concl_sz, _ = solve_card_fonts(
            [it['points'] for it in items], body_w, ch, SP_RANGE, DOMAINS,
            line_spacing=ls, max_lines=2)
        band = max(lh(num_sz), lh(title_sz))
        para_n = len(items[0]['points']) - 1
    else:                                   # 企业风固定节奏（28/16/12/16）
        body_sz, title_sz, num_sz, concl_sz = (
            vg.get('body_size', 12), vg.get('title_size', 16),
            vg.get('num_size', 28), vg.get('tag_size', 16))
        band = max(lh(num_sz), lh(title_sz))
        para_n = len(items[0]['points']) - 1
    # ── 四卡 ──
    positions = [(0, 0), (1, 0), (1, 1), (0, 1)]          # 左上→右上→右下→左下
    for it, (r, c) in zip(items, positions):
        x = x0 + c * (cw2 + gap_h)
        y = y0 + r * (ch + gap_v)
        content_card(slide, x, y, cw2, ch, line, line_w=line_w)
        if vg.get('solve', False):
            rows_k = [max(1, math.ceil(measure_text_width(p, body_sz) * 1.03 / body_w))
                      for p in it['points']]
            G = fit_gaps(sum(rows_k), body_sz, band, concl_sz, ch, SP_RANGE, para_n,
                         line_spacing=ls, line_thick=0.014)
            top, band_to_d1, d1_to_body = G['top'], G['band_to_d1'], G['d1_to_body']
            para, body_to_d2, d2_to_concl = G['para'], G['body_to_d2'], G['d2_to_concl']
        else:
            top, band_to_d1, d1_to_body = 0.20, 0.06, 0.06
            para, body_to_d2, d2_to_concl = 0.30, 0.08, 0.07
        # 带：编号 + 标题同行（垂直居中，行带高=max 行高，防数字压线）
        tf = tb_box(slide, x + PAD, y + top, 0.70, band + 0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, it['num'], num_sz, True, vg.get('num_color', cd['title_color']),
                     PP_ALIGN.LEFT, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = head_f; set_chinese_font(r2, head_f)
        tf = tb_box(slide, x + PAD + 0.75, y + top, cw2 - PAD - 0.75, band + 0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, it['title'], title_sz, True, vg.get('title_color', cd['title_color']),
                     PP_ALIGN.LEFT, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = head_f; set_chinese_font(r2, head_f)
        yy = y + top + band + band_to_d1
        # D1（可选：标题下青短线，不通栏）
        if vg.get('d1', False):
            head_w = 0.75 + measure_text_width(it['title'], title_sz) + 0.10
            card_divider(slide, x + PAD, yy, body_w + (0.16 if vg.get('mark', False) else 0.0),
                         line, pt=0.75, full=False, title_w=min(head_w, body_w))
            yy += 0.007
        # 要点：mark 方块（文质）或「· 」前缀（企业）
        rows_k = [max(1, math.ceil(measure_text_width(p, body_sz) * 1.03 / body_w))
                  for p in it['points']]
        for k, pt in enumerate(it['points']):
            if vg.get('mark', False):
                mark_y = yy + (lh(body_sz) - 0.09) / 2.0
                rect(slide, x + PAD, mark_y, 0.09, 0.09,
                     fill=vg.get('num_color', cd['title_color']), line=None)
                tf = tb_box(slide, x + PAD + 0.16, yy, body_w, rows_k[k] * lh(body_sz) + 0.04)
            else:
                tf = tb_box(slide, x + PAD, yy, body_w, rows_k[k] * lh(body_sz) + 0.04)
                p = run_text(tf, "· " + pt, body_sz, False, vg.get('body_color', cd['body_color']),
                             PP_ALIGN.LEFT, Pt(0), True)
                for r2 in p.runs:
                    r2.font.name = body_f; set_chinese_font(r2, body_f)
                yy += rows_k[k] * lh(body_sz) + (para if k < len(rows_k) - 1 else 0.0)
                continue
            p = run_text(tf, pt, body_sz, False, vg.get('body_color', cd['body_color']),
                         PP_ALIGN.LEFT, Pt(0), True)
            for r2 in p.runs:
                r2.font.name = body_f; set_chinese_font(r2, body_f)
            yy += rows_k[k] * lh(body_sz) + (para if k < len(rows_k) - 1 else 0.0)
        # D2（结论上分界，通栏）
        d2_y = yy + body_to_d2
        card_divider(slide, x + PAD, d2_y, body_w + (0.16 if vg.get('mark', False) else 0.0),
                     vg.get('div', cd['line']), pt=1.0, full=True)
        # 结论 tag（对齐由 tag_align 控制；居中时保持整宽 body_w，避免 fit_single_line 把框缩窄成伪左对齐）
        tf = tb_box(slide, x + PAD, d2_y + d2_to_concl, body_w, lh(concl_sz) + 0.04)
        tag_align = vg.get('tag_align', 'center')
        if tag_align == 'center':
            align = PP_ALIGN.CENTER          # 整宽内居中，结论在卡内容区真正左右居中
        else:
            fit_single_line(tf, it['tag'], concl_sz, buffer=0.2, max_w=body_w)  # 左对齐：收缩框宽防折行
            align = PP_ALIGN.LEFT
        p = run_text(tf, it['tag'], concl_sz, True, vg.get('tag_color', cd['concl_color']),
                     align, Pt(0), True)
        for r2 in p.runs:
            r2.font.name = body_f; set_chinese_font(r2, body_f)
    # ── 循环箭头（loop=True：4 段等长，四卡间距全等）──
    if loop:
        ar = skin['arrow']
        top_y = y0 + ch / 2
        bot_y = y0 + ch + gap_v + ch / 2
        l_mid = x0 + cw2 / 2
        r_mid = x0 + cw2 + gap_h + cw2 / 2
        loops = [
            (x0 + cw2, top_y, x0 + cw2 + gap_h, top_y),
            (r_mid, y0 + ch, r_mid, y0 + ch + gap_v),
            (x0 + cw2 + gap_h, bot_y, x0 + cw2, bot_y),
            (l_mid, y0 + ch + gap_v, l_mid, y0 + ch),
        ]
        for (x1, y1, x2, y2) in loops:
            loop_arrow(slide, x1, y1, x2, y2, color=ar['color'], w=ar['w'])
    return y_end


def layout_profile_warning(slide, skin, content, y0=1.85, y_end=6.30,
                           x0=CANVAS_X0, cw=CANVAS_CW):
    """画像+清醒提醒公共骨架（企业风 profile_warning 样板）：非对称双卡 + 结论行 + 两级分隔线。

    content = {left:{title, roles[], conclusion}, right:{title, warns[{kw,desc}], conclusion}}
    roles 短词(≤col_threshold 字)自动两列、长词独占行。
    skin['profile'] 字段：
      card_bg(实心卡底) · title_size/title_color · title_line(标题下细线色)
      role_size/role_color/dot(圆点色)/col_threshold
      warn_kw_size/warn_kw_color/warn_desc_size/warn_desc_color/mark(方块色)
      faint_div(行间浅线)/strong_div(结论上方亮线)
      concl_bar(结论竖线)/concl_size/concl_color
    页眉与底部金句由调用方处理。"""
    pf = skin.get('profile', {})
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    LW = 4.90                          # 左窄卡宽（企业风 0.60→5.50）
    RX = x0 + 5.30                     # 右卡 x（企业风 5.90）
    RW = cw + x0 - RX - 0.10           # 右宽（与内容区右缘对齐）
    PAD = 0.30
    VPAD = 0.30
    title_sz = pf.get('title_size', 18)
    threshold = pf.get('col_threshold', 8)
    kh_sz = pf.get('warn_kw_size', 16)
    kd_sz = pf.get('warn_desc_size', 12)
    concl_sz = pf.get('concl_size', 14)
    # ── 卡高自适应内容（框紧抱内容 + 整簇垂直居中），消死白 ──
    def _left_h(left):
        shorts = [r for r in left['roles'] if len(r) <= threshold]
        longs = [r for r in left['roles'] if len(r) > threshold]
        n_rows = math.ceil(len(shorts) / 2.0) + len(longs)
        roles_h = max(1, n_rows) * 0.40 + 0.10
        return 0.72 + 0.20 + roles_h + 0.30 + (lh(concl_sz) + 0.40)
    def _right_h(right):
        wh = 0.0
        for it in right['warns']:
            wh += lh(kh_sz) + 0.06 + lh(kd_sz) * 2 + 0.06
        return 0.72 + 0.30 + wh + 0.30 + (lh(concl_sz) + 0.40)
    CH = max(_left_h(content['left']), _right_h(content['right'])) + 2 * VPAD
    CH = min(CH, y_end - y0)
    CY = y0 + max(0.0, (y_end - y0 - CH) / 2.0)
    def _card(x, w, title):
        rounded_card(slide, x, CY, w, CH, fill=pf.get('card_bg'), adj=0.04)
        tf = tb_box(slide, x + PAD, CY + 0.28, w - PAD * 2, lh(title_sz) + 0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, title, title_sz, True, pf.get('title_color', WHITE),
                     PP_ALIGN.LEFT, Pt(0), True)
        for r in p.runs:
            r.font.name = head_f; set_chinese_font(r, head_f)
        rect(slide, x + PAD, CY + 0.72, measure_text_width(title, title_sz) + 0.30, Pt(1.5),
             fill=pf.get('title_line', skin['card']['line']))
    def _conclusion(x, w, text):
        DIV_Y = CY + CH - 0.86
        CONC_Y = CY + CH - 0.72
        rect(slide, x + PAD, DIV_Y, w - PAD * 2, Pt(1), fill=pf.get('strong_div'))
        rect(slide, x + PAD, CONC_Y + 0.03, Pt(2), 0.32, fill=pf.get('concl_bar'))
        tf = tb_box(slide, x + PAD + 0.15, CONC_Y, w - PAD * 2 - 0.15, 0.40)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, text, concl_sz, True, pf.get('concl_color'),
                     PP_ALIGN.LEFT, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
    def _role(x, y, w, text, size=None, color=None):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x),
            Inches(y + lh(size or pf.get('role_size', 13)) / 2 - 0.05), Inches(0.10), Inches(0.10))
        dot.fill.solid(); dot.fill.fore_color.rgb = pf.get('dot', skin['card']['line'])
        dot.line.fill.background(); dot.shadow.inherit = False
        tf = tb_box(slide, x + 0.23, y, w - 0.23, 0.34)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, text, size or pf.get('role_size', 13), False,
                     color or pf.get('role_color', LIGHT_GRAY), PP_ALIGN.LEFT, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
    # ── 左卡：人/角色罗列（短词两列 + 长词独占行），卡内均匀铺开 ──
    left = content['left']
    _card(x0, LW, left['title'])
    LI_X, LI_W = x0 + PAD, LW - PAD * 2
    COL_W = LI_W / 2.0
    shorts = [r for r in left['roles'] if len(r) <= threshold]
    longs = [r for r in left['roles'] if len(r) > threshold]
    short_rows = math.ceil(len(shorts) / 2.0)
    total_rows = short_rows + len(longs)
    region_top = CY + 1.00
    region_bot = CY + CH - 1.00
    step = (region_bot - region_top) / max(1, total_rows) if total_rows else 0
    row_idx = 0
    for j in range(short_rows):
        gy = region_top + (row_idx + 0.5) * step
        for gc in (0, 1):
            if j * 2 + gc < len(shorts):
                _role(LI_X + gc * COL_W, gy - 0.17, COL_W, shorts[j * 2 + gc])
        row_idx += 1
    for r in longs:
        gy = region_top + (row_idx + 0.5) * step
        _role(LI_X, gy - 0.17, LI_W, r)
        row_idx += 1
    _conclusion(x0, LW, left['conclusion'])
    # ── 右卡：观点/告诫（方块 + kw/desc），卡内均匀铺开 ──
    right = content['right']
    _card(RX, RW, right['title'])
    R_IN, R_INW = RX + PAD, RW - PAD * 2
    warns = right['warns']
    wregion_top = CY + 1.00
    wregion_bot = CY + CH - 1.00
    wstep = (wregion_bot - wregion_top) / max(1, len(warns)) if warns else 0
    for i, it in enumerate(warns):
        kw_y = (wregion_top + (i + 0.5) * wstep) - 0.38
        rect(slide, R_IN, kw_y + lh(kh_sz) / 2 - 0.04, 0.085, 0.085,
             fill=pf.get('mark', skin['card']['line']))
        _skin_text(slide, R_IN + 0.20, kw_y, R_INW - 0.20, lh(kh_sz) + 0.05,
                   it['kw'], kh_sz, True, pf.get('warn_kw_color'), body_f)
        _skin_text(slide, R_IN, kw_y + lh(kh_sz) + 0.06, R_INW, lh(kd_sz) * 2 + 0.05,
                   it['desc'], kd_sz, False, pf.get('warn_desc_color'), body_f)
    _conclusion(RX, RW, right['conclusion'])
    return y_end


def layout_six_step(slide, skin, content, y0=1.88, y_end=6.30,
                    x0=CANVAS_X0, cw=CANVAS_CW):
    """六步流程 + 三栏卡 + 底部说明公共骨架（文质六步流程样板，企业风节奏）。

    content = {cap, steps[6 str], cards[{title, lines[]}×3], note}
    步骤条：6 小卡 + 顶行箭头 + U 型回流（起终点步首/步尾底中，走 common.loop_return_u）；
    三栏卡；note 悬挂引注（皮肤 tail_bar）。"""
    cap, cd, ar = skin['cap'], skin['card'], skin['arrow']
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    y = y0
    if content.get('cap'):
        _skin_text(slide, x0, y, cw, 0.30, content['cap'], cap['size'], True, cap['color'], head_f)
        y += lh(cap['size']) + 0.12
    steps = content['steps']
    sw, sg = (cw - 5 * 0.22) / 6.0, 0.22
    sy, sh = y, 0.66
    centers = []
    for i, name in enumerate(steps):
        x = x0 + i * (sw + sg)
        content_card(slide, x, sy, sw, sh, cd['line'], line_w=cd.get('line_w', 1.25))
        _skin_text(slide, x + 0.04, sy + 0.08, sw - 0.08, 0.20, "0%d" % (i + 1),
                   cd['body_size'], True, cd['title_color'], body_f, PP_ALIGN.CENTER)
        _skin_text(slide, x + 0.04, sy + 0.32, sw - 0.08, 0.28, name,
                   cd['body_size'], True, cd['title_color'], body_f, PP_ALIGN.CENTER)
        centers.append(x + sw / 2.0)
        if i < len(steps) - 1:
            loop_arrow(slide, x + sw + 0.04, sy + sh / 2.0, x + sw + sg - 0.04, sy + sh / 2.0,
                       color=ar['color'], w=ar['w'])
    loop_return_u(slide, centers, sy + sh, sy + sh + 0.22, color=ar['color'], w=ar['w'])
    y = sy + sh + 0.22 + 0.24
    bw = (cw - 2 * 0.30) / 3.0
    ch = 1.88
    for i, c in enumerate(content['cards']):
        x = x0 + i * (bw + 0.30)
        content_card(slide, x, y, bw, ch, cd['line'], line_w=cd.get('line_w', 1.25))
        _skin_text(slide, x + 0.26, y + 0.22, bw - 0.52, 0.34, c['title'],
                   cd['title_size'], True, cd['title_color'], head_f)
        yy = y + 0.76
        for ln in c['lines']:
            h = max(1, math.ceil(measure_text_width(ln, cd['body_size']) * 1.03 / (bw - 0.52))) \
                * lh(cd['body_size'])
            _skin_text(slide, x + 0.26, yy, bw - 0.52, h + 0.04, ln,
                       cd['body_size'], False, cd['body_color'], body_f)
            yy += h + 0.12
    y += ch + 0.22
    if content.get('note'):
        bar = skin.get('tail_bar')
        tx, tw = (x0 + 0.16, cw - 0.16) if bar else (x0, cw)
        h = max(1, math.ceil(measure_text_width(content['note'], 14) * 1.03 / tw)) * 14 * ls / 72.0
        if bar:
            rect(slide, x0, y + 0.02, 0.022, h, fill=bar)
        _skin_text(slide, tx, y, tw, h + 0.08, content['note'], 14, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f)
    return y


def layout_talent_strip(slide, skin, content, y0=1.88, y_end=6.30,
                        x0=CANVAS_X0, cw=CANVAS_CW):
    """人才条公共骨架：intro + 全宽横条列表（num + title + lines）。

    content = {intro, items[{num, title, lines[]}×4]}；卡透明底+皮肤强调色边框。"""
    cd = skin['card']
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    y = y0
    if content.get('intro'):
        bar = skin.get('tail_bar')
        tx, tw = (x0 + 0.16, cw - 0.16) if bar else (x0, cw)
        h = max(1, math.ceil(measure_text_width(content['intro'], 14) * 1.03 / tw)) * 14 * ls / 72.0
        if bar:
            rect(slide, x0, y + 0.02, 0.022, h, fill=bar)
        _skin_text(slide, tx, y, tw, h + 0.08, content['intro'], 14, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f)
        y += h + 0.24
    sh, sg = 0.82, 0.04
    for i, it in enumerate(content['items']):
        yy = y + i * (sh + sg)
        content_card(slide, x0, yy, cw, sh, cd['line'], line_w=cd.get('line_w', 1.25))
        _skin_text(slide, x0 + 0.28, yy + 0.10, 0.55, 0.30, it['num'],
                   cd['title_size'], True, cd['title_color'], head_f)
        _skin_text(slide, x0 + 0.86, yy + 0.11, 5.0, 0.30, it['title'],
                   cd['title_size'], True, cd['title_color'], head_f)
        ly = yy + 0.40
        for ln in it['lines']:
            h = max(1, math.ceil(measure_text_width(ln, cd['body_size']) * 1.03 / (cw - 0.56))) \
                * lh(cd['body_size'])
            _skin_text(slide, x0 + 0.28, ly, cw - 0.56, h + 0.02, ln,
                       cd['body_size'], False, cd['body_color'], body_f)
            ly += h + 0.02
    return y + len(content['items']) * (sh + sg)


def layout_transition_rows(slide, skin, content, y0=1.88, y_end=6.30,
                           x0=CANVAS_X0, cw=CANVAS_CW):
    """转型对照公共骨架：每行 = 旧态降温框 → 箭头 → 新态青边高亮卡（企业风 transform_summary 样板）。

    content = {cap, rows[{from, to_title, to_lines[]}]}；行高随行数自适应。
    skin['transition'] = {from_bg(旧态底), from_color(旧态字), to_color(新态标题),
                          line(新态边/竖线)}，缺省回退 skin['card']。"""
    tr = skin.get('transition', {})
    cd = skin['card']
    cap = skin['cap']
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    n0 = len(slide.shapes)
    y = y0
    if content.get('cap'):
        _skin_text(slide, x0, y, cw, 0.30, content['cap'], cap['size'], True, cap['color'], head_f)
        y += lh(cap['size']) + 0.14
    rows = content['rows']
    rh = (y_end - y - 0.13 * (len(rows) - 1)) / len(rows)
    numbered = tr.get('numbered', False)
    if numbered:                                     # 企业风定稿：编号(青) + 旧态框
        num_w, from_x, from_w = 0.65, x0 + 0.65, 3.00
        arrow_w = 0.60
    else:
        num_w, from_x, from_w = 0.0, x0, 3.30
        arrow_w = 0.52
    to_x = from_x + from_w + arrow_w
    to_w = cw + x0 - to_x - 0.10
    ar = skin['arrow']
    for i, it in enumerate(rows):
        yy = y + i * (rh + 0.13)
        if numbered:
            tf = tb_box(slide, x0, yy, num_w - 0.10, rh)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = run_text(tf, "0%d" % (i + 1), 20, True, skin['cap']['color'],
                         PP_ALIGN.CENTER, Pt(0), True)
            for r in p.runs:
                r.font.name = body_f; set_chinese_font(r, body_f)
        rect(slide, from_x, yy, from_w, rh, fill=tr.get('from_bg'))
        tf = tb_box(slide, from_x + 0.16, yy, from_w - 0.32, rh)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, it['from'], 16, True, tr.get('from_color'), PP_ALIGN.CENTER, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
        loop_arrow(slide, from_x + from_w + 0.10, yy + rh / 2.0, to_x - 0.02, yy + rh / 2.0,
                   color=ar['color'], w=ar['w'])
        content_card(slide, to_x, yy, to_w, rh, tr.get('line', cd['line']),
                     line_w=cd.get('line_w', 1.25))
        rect(slide, to_x, yy + 0.18, 0.022, rh - 0.36, fill=tr.get('line', cd['line']))
        # 新态卡内：标题 + 描述 在 rh 内垂直居中；标题↔描述间距 ≥0.10（对齐 SKILL §4.12-7）
        _tsz = 16
        _dsz = cd['body_size']
        _dl = it.get('to_lines', [])
        _dh = sum(max(1, math.ceil(measure_text_width(ln, _dsz) * 1.03 / (to_w - 0.52))) * lh(_dsz)
                  for ln in _dl)
        _gap = 0.12
        _bh = 0.30 + _gap + _dh
        _bt = yy + max(0.0, (rh - _bh) / 2.0)
        _skin_text(slide, to_x + 0.26, _bt, to_w - 0.52, 0.30, it['to_title'],
                   _tsz, True, tr.get('to_color', cd['title_color']), head_f)
        ly = _bt + 0.30 + _gap
        for ln in _dl:
            h = max(1, math.ceil(measure_text_width(ln, _dsz) * 1.03 / (to_w - 0.52))) * lh(_dsz)
            _skin_text(slide, to_x + 0.26, ly, to_w - 0.52, h + 0.02, ln,
                       _dsz, False, cd['body_color'], body_f)
            ly += h + 0.02
    yb = y + len(rows) * (rh + 0.13)
    _vshift(slide, n0, y0, yb, y_end)
    return yb


def layout_hero_questions(slide, skin, content, y0=1.88, y_end=6.30,
                          x0=CANVAS_X0, cw=CANVAS_CW):
    """结语 hero 问答公共骨架：lead + questions 大字居中 + 强调短线 + signature。
    二维码属氛围装饰留引擎层（调用方自行 add_qr）。
    content = {lead, questions, signature}；questions 字号 skin['hero']['size']（缺省 36）。"""
    cd = skin['card']
    body_f = skin['fonts']['body']
    qsz = skin.get('hero', {}).get('size', 36)
    ls = skin.get('line_spacing', 1.2)
    y = y0 + 0.40
    if content.get('lead'):
        _skin_text(slide, x0, y, cw, 0.34, content['lead'], 18, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f, PP_ALIGN.CENTER)
        y += 0.44
    # 实测问题句折行数，按真实高度定框（避免 2 行时文字与下方重叠）
    qw = measure_text_width(content['questions'], qsz)
    qrows = max(1, math.ceil(qw * 1.03 / cw))
    qh = qrows * (qsz * ls / 72.0)
    tf = tb_box(slide, x0, y, cw, qh + 0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = run_text(tf, content['questions'], qsz, True, cd['title_color'],
                 PP_ALIGN.CENTER, Pt(0), True)
    for r in p.runs:
        r.font.name = skin['fonts']['head']; set_chinese_font(r, skin['fonts']['head'])
    # Q4：去掉强调短线（用户拍板），签名直接落在真实文字底下方
    y += qh + 0.45
    if content.get('signature'):
        _skin_text(slide, x0, y, cw, 0.30, content['signature'], 14, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f, PP_ALIGN.CENTER)
    return y


def layout_two_column(slide, skin, content, y0=1.88, y_end=6.30,
                      x0=CANVAS_X0, cw=CANVAS_CW):
    """双栏并列公共骨架：两卡（title + items），等体裁等体量。
    content = {columns[{title, items[]}×2]}。"""
    cd = skin['card']
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    ls = skin.get('line_spacing', 1.2)
    def lh(sz): return sz * ls / 72.0
    cw2 = (cw - 0.67) / 2.0
    ch = y_end - y0
    for i, col in enumerate(content['columns']):
        x = x0 + i * (cw2 + 0.67)
        content_card(slide, x, y0, cw2, ch, cd['line'], line_w=cd.get('line_w', 1.25))
        _skin_text(slide, x + 0.35, y0 + 0.30, cw2 - 0.70, 0.40, col['title'],
                   cd['title_size'], True, cd['title_color'], head_f)
        yy = y0 + 0.95
        for it in col['items']:
            h = max(1, math.ceil(measure_text_width("·  " + it, cd['body_size']) * 1.03 / (cw2 - 0.70))) \
                * lh(cd['body_size'])
            _skin_text(slide, x + 0.35, yy, cw2 - 0.70, h + 0.04, "·  " + it,
                       cd['body_size'], False, cd['body_color'], body_f)
            yy += h + 0.16
    return y_end


def layout_summary(slide, skin, content, y0=1.88, y_end=6.30,
                   x0=CANVAS_X0, cw=CANVAS_CW):
    """总结页公共骨架：title 大字居中 + sub + metrics 行 + quote 引文。
    content = {title, sub, metrics[], quote}。"""
    cd = skin['card']
    head_f, body_f = skin['fonts']['head'], skin['fonts']['body']
    y = y0 + 0.2
    tf = tb_box(slide, x0, y, cw, 1.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = run_text(tf, content['title'], 46, True, cd['title_color'], PP_ALIGN.CENTER, Pt(12), True)
    for r in p.runs:
        r.font.name = head_f; set_chinese_font(r, head_f)
    y += 1.2
    if content.get('sub'):
        _skin_text(slide, x0, y, cw, 0.40, content['sub'], 20, False,
                   skin['fonts'].get('mute', cd['body_color']), body_f, PP_ALIGN.CENTER)
        y += 0.6
    if content.get('metrics'):
        tf = tb_box(slide, x0, y, cw, 0.6)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        for m in content['metrics']:
            r = p.add_run(); r.text = m
            r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = cd['concl_color']
            r.font.name = body_f; set_chinese_font(r, body_f)
            rr = p.add_run(); rr.text = "    "
            rr.font.size = Pt(14); rr.font.color.rgb = cd['body_color']; rr.font.name = body_f
        y += 0.8
    if content.get('quote'):
        rect(slide, (13.333 - 7.0) / 2.0, y, 7.0, 0.02, fill=cd['line'])
        tf = tb_box(slide, 1.8, y + 0.2, 9.73, 0.9)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = run_text(tf, content['quote'], 18, True, cd['concl_color'], PP_ALIGN.CENTER, Pt(0), True)
        for r in p.runs:
            r.font.name = body_f; set_chinese_font(r, body_f)
    return y
