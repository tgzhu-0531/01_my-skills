# -*- coding: utf-8 -*-
"""
fengmang.py — 锋芒风引擎（纯黑高对比 / Neo-Swiss 大字报）
黑底 #0A0A0A + 品牌青 #00A7CB + 巨号粗体 + 1px 暗青边框卡片。
暴露统一版式原语：cover / section_header / timeline / two_column / card_grid / summary / bottom_gold
调用时由 per-article runner 把 content 映射进这些原语。
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from common import (BIZ_TEAL, WHITE, LIGHT_GRAY, GOLD,
                    SP_UNIT, SP_GROUP, SP_CARD_BODY, SP_RUN_DEF,
                    tb_box, run_text, rect, add_qr, add_page_number, set_bg,
                    blank_deck, set_chinese_font,
                    GOLD_LINE_Y, GOLD_TEXT_Y,
                    bottom_gold as _gold, tb_box_hug, shape_bottom_in, cover_qr,
                    cover_line, fit_cover_title,
                    layout_timeline, layout_card_grid, layout_three_column, layout_two_column, layout_summary,
                    layout_compare, layout_loop_page, layout_profile_warning,
                    layout_transition_rows)

BG    = RGBColor(0x0A, 0x0A, 0x0A)
MUTE  = RGBColor(0x99, 0x99, 0x99)   # 次要文字（≈白 0.6α）
MUTE2 = RGBColor(0x66, 0x66, 0x66)   # 更弱
# 暗底封面灰阶（SKILL §4.5：纯黑底副标题≥#CCC≈12:1、辅助≥#B3B3B3≈9:1，杜绝"看不清"）
COVER_SUB = RGBColor(0xCC, 0xCC, 0xCC)   # 封面副标题（高亮档）
COVER_TAG = RGBColor(0xB3, 0xB3, 0xB3)   # 封面标签行（高亮档）
LINE_DIM = RGBColor(0x2A, 0x2A, 0x2A)
CARD_BG   = RGBColor(0x11, 0x11, 0x11)   # 锋芒风卡片底
BORDER_DIM = RGBColor(0x1A, 0x4A, 0x55)  # 锋芒风暗青边框（模拟 0.2α）
FENG_CARD_LINE = RGBColor(0x44, 0x44, 0x44)  # 内容卡边框（#444 对 #0A0A0A 底对比度≈8:1，清晰可见且不抢戏）

PAGE_W = 13.333

def _alpha(fill, a):
    """设置纯色填充透明度 a(0..1)。"""
    srgb = fill.fore_color._xFill.find(qn('a:srgbClr'))
    if srgb is None:
        return
    a_el = srgb.find(qn('a:alpha'))
    if a_el is None:
        a_el = etree.SubElement(srgb, qn('a:alpha'))
    a_el.set('val', str(int(a * 1000)))

def _center_x(w):
    return (PAGE_W - w) / 2.0

# ═══════════════════════════════════════════
# 画布 / 页管理（纯代码绘制，不克隆模板；与 business/enterprise 同签名）
# ═════════════════════════════════════════
def new_deck(template=None, total=30):
    """创建空白 16:9 画布（纯黑底，不依赖企业模板）。"""
    return blank_deck()

def first_slide(prs):
    """新增空白页作封面载体（纯代码绘制，不复用模板占位符）。"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, BG)
    return s

def add_content(prs):
    """新增内容页（空白版式，各原语内部 set_bg 上纯黑底）。"""
    return prs.slides.add_slide(prs.slide_layouts[6])

# ═══════════════════════════════════════════
# 版式原语
# ═══════════════════════════════════════════
def cover(slide, title, subtitle, tags=None, total=1, n=1):
    set_bg(slide, BG)
    cx = PAGE_W / 2.0
    # 主标题：居中大标题（锋芒封面 DNA=居中）；字号智能算，占屏≤72%
    title_pt = fit_cover_title(title, 54, min_pt=28)
    tf = cover_line(slide, 0.8, 1.05, 11.73, title, title_pt, WHITE, True,
                    PP_ALIGN.CENTER)
    # 青色分隔线（居中短条）：y 由标题实际框底推导，杜绝硬编码碰撞
    y = shape_bottom_in(tf) + 0.14
    bar_w = 1.6
    rect(slide, cx - bar_w / 2, y, bar_w, 0.06, fill=BIZ_TEAL, line=None)
    # 副标题（居中，暗底高亮档 #CCC）
    y += 0.06 + 0.22
    tf = cover_line(slide, 0.8, y, 11.73, subtitle, 26, COVER_SUB, False,
                    PP_ALIGN.CENTER)
    # 标签行（居中，#B3B3B3）
    if tags:
        y = shape_bottom_in(tf) + 0.16
        cover_line(slide, 0.8, y, 11.73, "   ·   ".join(tags), 14, COVER_TAG, False,
                   PP_ALIGN.CENTER)
    # 封面二维码（SKILL §4.9 底中）
    cover_qr(slide, "fengmang")
    add_page_number(slide, n, total, MUTE2)

def section_header(slide, h1, label, sub, n=None, total=None):
    set_bg(slide, BG)
    # 主标题：上移(0.7→0.35) + 缩小(40→26pt) + 强制 MIDDLE（包住+居中，根治"距顶远/过大/不居中"）
    title_tf = cover_line(slide, 0.8, 0.35, 11.73, h1, 26, WHITE, True, PP_ALIGN.LEFT)
    # 青色分隔线：由标题实际底推导，杜绝硬编码重叠（原 1.62 贴脸标题）
    line_y = shape_bottom_in(title_tf) + 0.12
    rect(slide, 0.8, line_y, 1.1, 0.06, fill=BIZ_TEAL, line=None)
    # 副标：线下方留紧凑呼吸 + MIDDLE 居中（label 青粗 + sub 灰，同一行内联，杜绝拆两行）
    sub_y = line_y + 0.06 + 0.06
    tf = tb_box(slide, 0.8, sub_y, 11.73, 0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = run_text(tf, label, 16, True, BIZ_TEAL, PP_ALIGN.LEFT, Pt(0), first=True)
    if sub:
        r2 = p.add_run(); r2.text = "    " + sub
        r2.font.size = Pt(16); r2.font.bold = False; r2.font.color.rgb = MUTE
        r2.font.name = '微软雅黑'; set_chinese_font(r2)
    if n and total:
        add_page_number(slide, n, total, MUTE2)
    return shape_bottom_in(tf)   # 返回副标底，供 timeline 推导 y0 防重叠

def timeline(slide, items, n=None, total=None, gold=None, y0=2.10, gold_kw_a=()):
    """垂直时间轴：公共骨架 layout_timeline（纯黑底白日期，皮肤 FENGMANG_SKIN.timeline）。
    y0 默认 2.10：由 section_header 副标底(≈1.88)+呼吸推导，根治'副标与时间线重叠'；
    y_end=6.05 与金句线(6.35)留 0.30″ 间隙。"""
    set_bg(slide, BG)
    layout_timeline(slide, FENGMANG_SKIN,
                    [{"date": it["time"], "text": it["text"],
                      "kw_teal": it.get("kw_teal", []), "kw_amber": it.get("kw_amber", []),
                      "amt": it.get("amt")}
                     for it in items],
                    y0=y0, y_end=6.05)
    if gold:
        _gold(slide, gold, color=BIZ_TEAL, size=16, kw_a=gold_kw_a)
    if n and total:
        add_page_number(slide, n, total, MUTE2)
def two_column(slide, left, right, n=None, total=None, h=4.1, gold=None):
    """双栏：公共骨架 layout_two_column（透明底+暗青边框）。"""
    set_bg(slide, BG)
    ye = 6.10 if gold else 6.30
    layout_two_column(slide, FENGMANG_SKIN, {"columns": [left, right]}, y0=2.30, y_end=ye)
    if gold:
        _gold(slide, gold, color=BIZ_TEAL, size=16)
    if n and total:
        add_page_number(slide, n, total, MUTE2)
def card_grid(slide, cards, cols=2, n=None, total=None):
    """卡片网格：公共骨架 layout_card_grid（透明底+暗青边框）。"""
    set_bg(slide, BG)
    layout_card_grid(slide, FENGMANG_SKIN, cards, cols=cols, y0=1.65, y_end=6.38)
    if n and total:
        add_page_number(slide, n, total, MUTE2)
def three_column(slide, title=None, label=None, cards=None, n=None, total=None, connect=False):
    """三列（锋芒风，委托通用 layout_three_column）：单排三卡、等宽等框高，与 2×2 一致（无下划线/无 icon/无连接线）。
    标题区复用 section_header（与全站一致）；n/total 由 section_header 绘制页码。"""
    set_bg(slide, BG)
    if title:
        section_header(slide, title, label, None, n=n, total=total)
    else:
        if n and total:
            add_page_number(slide, n, total, MUTE2)
    layout_three_column(slide, FENGMANG_SKIN, cards, y0=1.65, y_end=6.45, connect=connect)
def summary(slide, title, sub, metrics=None, quote=None, n=None, total=None):
    """收尾页：公共骨架 layout_summary（居中大标题+指标+金句）。"""
    set_bg(slide, BG)
    layout_summary(slide, FENGMANG_SKIN, {"title": title, "sub": sub,
                                          "metrics": metrics or [], "quote": quote},
                   y0=1.70, y_end=6.30)
    if n and total:
        add_page_number(slide, n, total, MUTE2)
def bottom_gold(slide, text, n=None, total=None, kw_a=()):
    """底部金句（青色居中，无分隔线）。机制委托 common.bottom_gold（共享金句坐标常量）。"""
    _gold(slide, text, color=BIZ_TEAL, size=16,
         text_y=GOLD_TEXT_Y, page_color=MUTE2,
         n=n, total=total, kw_a=kw_a)

# ═══════════════════════════════════════════
# 富内容原语（委托 common.py 内容专用骨架 + FENGMANG_SKIN；签名对齐 enterprise，
# 自含 section_header / 底部金句。纯黑高对比 Neo-Swiss 调性。）
# ═════════════════════════════════════════
def definition_compare(slide, h1, label, hero, hero_en, hero_cn,
                       left_bullets, left_summary, compare_header,
                       negatives, positive, gold=None, n=None, total=None):
    """定义+排除对比：公共骨架 layout_compare（纯黑底：白 hero / 青强调 / 暗灰圆底×badge）。"""
    set_bg(slide, BG)
    section_header(slide, h1, label, "")
    layout_compare(slide, FENGMANG_SKIN, {
        "hero": hero, "hero_en": hero_en, "hero_cn": hero_cn,
        "bullets": left_bullets, "summary": left_summary,
        "right": {"header": compare_header,
                  "negatives": [{"kw": kw, "desc": d} for kw, d in negatives],
                  "positive": {"kw": "", "desc": positive}},
    }, y0=1.88, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, MUTE2)

def work_people(slide, h1, label, loop_nodes, roles, abilities,
                gold=None, n=None, total=None):
    """「如何工作，谁适合做？」：公共骨架 layout_loop_page（①闭环 ②双卡⇄ ③能力卡）。"""
    set_bg(slide, BG)
    section_header(slide, h1, label, "")
    content = {
        "cap1": "① 工作方式：业务现场的一体化闭环",
        "nodes": [nd if isinstance(nd, dict) else {"name": nd} for nd in loop_nodes],
        "cap2": "② 两类角色：边听边做，边做边改",
        "cards": [{"title": it["title"], "desc": it["desc"]} for it in roles],
        "cap3": "③ 三项核心能力",
        "tail": [{"title": it["title"], "desc": it["desc"]} for it in abilities],
    }
    layout_loop_page(slide, FENGMANG_SKIN, content, y0=1.88, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, MUTE2)

def profile_warning(slide, h1, label, left, right, gold=None, n=None, total=None):
    """「人才画像 + 清醒提醒」：公共骨架 layout_profile_warning（非对称双卡+结论）。"""
    set_bg(slide, BG)
    section_header(slide, h1, label, "")
    layout_profile_warning(slide, FENGMANG_SKIN, {"left": left, "right": right},
                           y0=1.88, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, MUTE2)

def transform_summary(slide, h1, label, rows, gold=None, n=None, total=None):
    """「转型对照」：公共骨架 layout_transition_rows（旧态→新态，行高自适应）。"""
    set_bg(slide, BG)
    section_header(slide, h1, label, "")
    layout_transition_rows(slide, FENGMANG_SKIN,
                           {"cap": "三个转变",
                            "rows": [{"from": it["from"], "to_title": it["to_title"],
                                      "to_lines": [it["to_desc"]]} for it in rows]},
                           y0=1.77, y_end=6.10 if gold else 6.30)
    if gold:
        bottom_gold(slide, gold, n=n, total=total)
    elif n and total:
        add_page_number(slide, n, total, MUTE2)


# ═══════════════════════════════════════════
# 锋芒风皮肤（布局骨架的视觉参数，2026-07-31 架构：布局在 common.layout_*，皮肤在本层）
# ═══════════════════════════════════════════
FENGMANG_SKIN = {
    "bg": BG,
    "fonts": {"head": "微软雅黑", "body": "微软雅黑", "mute": MUTE},
    "cap":  {"size": 14, "bold": True, "color": BIZ_TEAL},
    "node": {"size": 16, "bold": True, "color": WHITE, "sub_size": 12, "sub_color": MUTE},
    "card": {"title_size": 16, "title_color": WHITE, "sub_size": 12, "sub_color": MUTE,
             "body_size": 12, "body_color": LIGHT_GRAY, "concl_color": BIZ_TEAL,
             "line": FENG_CARD_LINE, "line_w": 1.0},
    "arrow": {"color": WHITE, "w": Pt(2)},       # 纯黑底 → 白 2pt
    "timeline": {"axis": BIZ_TEAL, "dot": BIZ_TEAL, "date_size": 14, "date_color": WHITE,
                 "event_size": 14, "event_color": LIGHT_GRAY, "arrow": False,
                 "kw_teal_color": BIZ_TEAL, "kw_amber_color": GOLD},
    # 定义对比（纯黑底：白 hero / 青强调 / 暗灰圆底×）
    "compare": {
        "hero_size": 48, "hero_color": WHITE,
        "en_size": 16, "en_color": LIGHT_GRAY, "en_italic": True,
        "cn_size": 18, "cn_color": BIZ_TEAL,
        "bullet_size": 14, "bullet_color": LIGHT_GRAY, "bullet_mark": "dot", "mark_color": BIZ_TEAL,
        "summary_size": 15, "summary_color": BIZ_TEAL,
        "header_size": 18, "header_color": WHITE,
        "neg_badge": "circle", "neg_kw_size": 18, "neg_kw_color": LIGHT_GRAY,
        "neg_desc_size": 15, "neg_desc_color": MUTE,
        "neg_badge_bg": RGBColor(0x1E, 0x1E, 0x1E), "neg_badge_color": MUTE,
        "pos_bg": RGBColor(0x14, 0x14, 0x14), "pos_border": FENG_CARD_LINE,
        "pos_kw_size": 14, "pos_kw_color": BIZ_TEAL,
        "pos_desc_size": 14, "pos_desc_color": LIGHT_GRAY,
        "div": BORDER_DIM, "right_card": False,
    },
    # 价值网格（纯黑底：白卡字 / 青锚点 / 等长循环箭头）
    "value_grid": {
        "num_size": 28, "num_color": BIZ_TEAL, "title_size": 16, "title_color": WHITE,
        "body_size": 12, "body_color": LIGHT_GRAY, "tag_size": 16, "tag_color": BIZ_TEAL,
        "tag_align": "center", "div": RGBColor(0x2A, 0x3A, 0x4D), "d1": False, "mark": False,
        "gap_h": 0.70, "gap_v": 0.70, "solve": False,
        "line": FENG_CARD_LINE, "line_w": 1.0,
    },
    # 画像（近黑实心卡 + 青结论竖线）
    "profile": {
        "card_bg": RGBColor(0x12, 0x12, 0x12), "title_size": 18, "title_color": WHITE,
        "title_line": BIZ_TEAL, "role_size": 13, "role_color": LIGHT_GRAY,
        "dot": BIZ_TEAL, "col_threshold": 8,
        "warn_kw_size": 16, "warn_kw_color": BIZ_TEAL,
        "warn_desc_size": 12, "warn_desc_color": LIGHT_GRAY, "mark": BIZ_TEAL,
        "faint_div": RGBColor(0x22, 0x22, 0x22), "strong_div": RGBColor(0x33, 0x33, 0x33),
        "concl_bar": BIZ_TEAL, "concl_size": 14, "concl_color": BIZ_TEAL,
    },
    # 转型对照（旧态近黑降温 / 新态青边高亮）
    "transition": {
        "from_bg": RGBColor(0x1E, 0x1E, 0x1E), "from_color": MUTE,
        "to_color": WHITE, "line": FENG_CARD_LINE,
    },
    "tail_bar": None,
    "line_spacing": 1.2,
}

def save(prs, outpath):
    """保存（纯代码绘制，无模板克隆，无需 noFill 清理）。"""
    prs.save(outpath)
    print("saved:", outpath, "slides:", len(prs.slides._sldIdLst))
