# -*- coding: utf-8 -*-
"""
enterprise.py — 企业风引擎（原版 biz_engine.py 直译）
模板克隆（保留品牌区） + 内容布局 + 企业配色。
与 tgzhu-content-2pptx/references/mck-ppt/biz_engine.py 完全一致，
仅从类方法改为函数式原语（cover / section_header / two_column / card_grid / timeline / bottom_gold）。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree
from common import (BIZ_TEAL, WHITE, LIGHT_GRAY, GOLD, ORANGE,
                    SP_UNIT, SP_GROUP, SP_CARD_BODY, SP_RUN_DEF,
                    tb_box, run_text, rect, add_qr, add_page_number,
                    set_chinese_font, remove_placeholders, _clean_shape, template_path, qr_path, FONT,
                    add_homeplate, add_dashed_line, fmt_date_5digit,
                    TYPOGRAPHY, apply_role, measure_text_width, find_cjk_font,
                    loop_arrow, LOOP_ARROW_W, loop_return_u, fit_single_line, fit_font_to_width,
                    cover_line, fit_cover_title, shape_bottom_in, rounded_card, emph_runs, add_soft_shadow,
                    bottom_gold as _gold, layout_loop_page, layout_card_grid, layout_two_column,
                    layout_three_column,
                    layout_compare, layout_value_grid, layout_profile_warning,
                    layout_transition_rows, content_card)

# ═══════════════════════════════════════════
# 原版 biz_engine 常量（全量搬运）
# ═══════════════════════════════════════════
BIZ_CARD_BG     = RGBColor(0x1A, 0x2D, 0x42)   # 深蓝卡底
BIZ_CARD_BORDER = RGBColor(0x00, 0xA7, 0xCB)    # 品牌青边框
BIZ_ACCENT      = RGBColor(0x00, 0xAF, 0xD2)    # 强调青（副标题）
ENTER_CARD_LINE  = RGBColor(0xD2, 0xD8, 0xDC)    # 内容卡极淡边框（贴浅底，比金句线稍可见）
ENT_MUTE        = RGBColor(0xCC, 0xCC, 0xCC)    # 页码/署名灰

# ═══════════════════════════════════════════
# 企业风皮肤（布局骨架的视觉参数，2026-07-31 架构：布局在 common.layout_*，皮肤在本层）
# ═══════════════════════════════════════════
ENTERPRISE_SKIN = {
    "bg": None,                                  # 深蓝渐变由母版承载
    "fonts": {"head": "微软雅黑", "body": "微软雅黑", "mute": ENT_MUTE},
    "cap":  {"size": 14, "bold": True, "color": BIZ_ACCENT},
    "node": {"size": 16, "bold": True, "color": WHITE, "sub_size": 12, "sub_color": LIGHT_GRAY},
    "card": {"title_size": 16, "title_color": WHITE, "sub_size": 12, "sub_color": LIGHT_GRAY,
             "body_size": 12, "body_color": LIGHT_GRAY, "concl_color": BIZ_ACCENT,
             "accent": BIZ_TEAL,
             "line": ENTER_CARD_LINE, "line_w": 1.25,
             # 双栏卡边框：近背景暗蓝，比原#D2D8DC浅灰更退背景层、不抢戏（仅 enterprise 双栏页生效）
             "two_col_border": RGBColor(0x22, 0x40, 0x5F),
             # 卡内分隔线专用暗色：退背景层（规范 §视觉层级），不与卡边框(浅灰)混淆
             "div_line": RGBColor(0x44, 0x58, 0x6C)},
    "arrow": {"color": WHITE, "w": Pt(2)},       # 深底 → 白 2pt
    "timeline": {"axis": BIZ_TEAL, "dot": BIZ_TEAL, "date_size": 14, "date_color": BIZ_ACCENT,
                 "event_size": 14, "event_color": LIGHT_GRAY, "arrow": True},
    # 定义对比（definition_compare 定稿：左宽右窄+竖线；neg_badge=circle 圆底×）
    "compare": {
        "hero_size": 48, "hero_color": BIZ_ACCENT,
        "en_size": 16, "en_color": WHITE, "en_italic": True,
        "cn_size": 18, "cn_color": BIZ_ACCENT,
        "bullet_size": 14, "bullet_color": LIGHT_GRAY, "bullet_mark": "dot", "mark_color": BIZ_ACCENT,
        "summary_size": 15, "summary_color": BIZ_ACCENT,
        "header_size": 18, "header_color": WHITE,
        "neg_badge": "circle", "neg_kw_size": 18, "neg_kw_color": RGBColor(0x9A, 0xA5, 0xB0),
        "neg_desc_size": 15, "neg_desc_color": RGBColor(0x7C, 0x88, 0x94),
        "neg_badge_bg": RGBColor(0x3A, 0x46, 0x54), "neg_badge_color": RGBColor(0x9A, 0xA5, 0xB0),
        "pos_bg": RGBColor(0x10, 0x33, 0x3E), "pos_border": ENTER_CARD_LINE,
        "pos_kw_size": 14, "pos_kw_color": BIZ_ACCENT,
        "pos_desc_size": 14, "pos_desc_color": BIZ_ACCENT,
        "div": RGBColor(0x2A, 0x3A, 0x4D), "right_card": False,
    },
    # 价值网格（value_loop 定稿：四卡等距 0.7 + 等长循环箭头；固定字号节奏）
    "value_grid": {
        "num_size": 28, "num_color": BIZ_ACCENT, "title_size": 16, "title_color": WHITE,
        "body_size": 12, "body_color": LIGHT_GRAY, "tag_size": 16, "tag_color": BIZ_ACCENT,
        "tag_align": "center", "div": RGBColor(0x44, 0x58, 0x6C), "d1": False, "mark": False,
        "gap_h": 0.50, "gap_v": 0.50, "solve": True,
        "card_w": "measured", "center": True,
        "line": ENTER_CARD_LINE, "line_w": 1.25,
    },
    # 画像（profile_warning 定稿：非对称实心卡 + 两级分隔线 + 青结论竖线）
    "profile": {
        "card_bg": RGBColor(0x16, 0x29, 0x3C), "title_size": 18, "title_color": WHITE,
        "title_line": BIZ_ACCENT, "role_size": 13, "role_color": LIGHT_GRAY,
        "dot": BIZ_ACCENT, "col_threshold": 8,
        "warn_kw_size": 16, "warn_kw_color": BIZ_ACCENT,
        "warn_desc_size": 12, "warn_desc_color": LIGHT_GRAY, "mark": BIZ_ACCENT,
        "faint_div": RGBColor(0x24, 0x34, 0x47), "strong_div": RGBColor(0x44, 0x58, 0x6C),
        "concl_bar": BIZ_ACCENT, "concl_size": 14, "concl_color": BIZ_ACCENT,
    },
    # 转型对照（transform_summary：旧态降温 #3A4654 / 新态青边高亮）
    "transition": {
        "from_bg": RGBColor(0x3A, 0x46, 0x54), "from_color": RGBColor(0x9A, 0xA5, 0xB0),
        "to_color": WHITE, "line": ENTER_CARD_LINE, "numbered": True,
    },
    "tail_bar": None,
    "line_spacing": 1.2,
}
ENT_DIVIDER     = RGBColor(0x1F, 0x2A, 0x38)    # 双栏 gutter 分隔线

BODY_SIZE   = Pt(12)   # 正文 12pt
SMALL_SIZE  = Pt(12)
FOOTNOTE    = Pt(9)
PAGE_W = 13.333

# ── 时间轴定稿版布局常量（沉淀自时间轴迭代验收）──
TL_SIDE      = 1.45     # 左右留空（对等）
TL_AXIS_X    = 1.45
TL_DATE_X    = 1.60
TL_DATE_W    = 1.05
TL_CARD_X    = 2.85
TL_DASH_X    = 2.75     # TL_CARD_X - 0.10
TL_CARD_W    = 9.00
TL_CARD_RIGHT= 11.85     # TL_CARD_X + TL_CARD_W（右留空 1.48 ≈ 左 1.45）
TL_CARD_ADJ  = 15       # homePlate 切角比例（小，微妙不三角）
TL_PAD_X     = 0.30
TL_PAD_Y     = 0.10
TL_SPACING   = 1.4
TL_AREA_TOP  = 1.85
TL_AREA_BOT  = 6.30
TL_DATE_SZ   = TYPOGRAPHY["date"][0]        # 20
TL_BODY_SZ   = TYPOGRAPHY["body"][0]        # 14
TL_BADGE_SZ  = TYPOGRAPHY["badge"][0]       # 12
TL_CARD_BORDER = RGBColor(0x1A, 0x55, 0x60)   # 暗青边框（接背景，不刺眼）
TL_BADGE_BG     = RGBColor(0x0F, 0x2A, 0x38)   # 徽章底
TL_DASH_CLR     = RGBColor(0x33, 0x66, 0x80)   # 每卡短虚线
TL_CYAN_DATE    = RGBColor(0x66, 0xC5, 0xDC)   # 日期亮青蓝

# ── 企业风角色配色（与 common.TYPOGRAPHY 角色一一对应）──
# 视觉语言：白(主) / 青(副) / 橙(结论)。金 #C9A227 在企业风退役。
ENT_PALETTE = {
    "hero":       WHITE,            # 封面主标题
    "h1":         WHITE,            # 页面/章节主标题
    "title_card": WHITE,            # 卡片标题
    "label":      BIZ_ACCENT,       # 副标题核心观点 = 强调青 #00AFD2
    "gold":       ORANGE,           # 底部金句 = 警示橙 #F79646
    "date":       TL_CYAN_DATE,     # 时间轴日期 = 亮青蓝 #66C5DC
    "body":       LIGHT_GRAY,       # 正文
    "body_small": LIGHT_GRAY,       # 密排卡正文
    "badge":      BIZ_ACCENT,       # 徽章
}

def _center_x(w):
    return (PAGE_W - w) / 2.0

# ═══════════════════════════════════════════
# 原版 add_text（逐行照抄 biz_engine.py line 80-101）
# ═══════════════════════════════════════════
def _add_text(slide, left, top, width, height, text,
              font_size=BODY_SIZE, font_name='微软雅黑',
              font_color=LIGHT_GRAY, bold=False,
              alignment=PP_ALIGN.LEFT, ea_font='微软雅黑',
              anchor=MSO_ANCHOR.TOP, line_spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True; tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns', 'tIns', 'rIns', 'bIns']:
        bodyPr.set(attr, '45720')
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = font_size; p.font.name = font_name
        p.font.color.rgb = font_color; p.font.bold = bold; p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0); p.space_after = Pt(0)
        p.line_spacing = 0.93 if font_size.pt >= 18 else Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_chinese_font(run, ea_font)
    return txBox

def _add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None: sp.remove(style)
    return shape

def _action_title(slide, text):
    """内容页标题（突出重点）：28pt Bold 白，x=0.46 y=0.25 w=8.6。与 section_header 的 H1(TYPOGRAPHY['h1']=28pt) 一致。"""
    _add_text(slide, Inches(0.46), Inches(0.25), Inches(8.6), Inches(0.65), text,
              font_size=Pt(28), font_color=WHITE, bold=True,
              font_name='微软雅黑', anchor=MSO_ANCHOR.BOTTOM)

# ═══════════════════════════════════════════
# 画布 / 页管理
# ═══════════════════════════════════════════
def new_deck():
    return Presentation(template_path())

def first_slide(prs):
    """模板自带第 1 页作封面载体（不清除任何形状，由 cover() 处理）。

    ⚠️ 铁律：克隆企业模板后**严禁删除 sldIdLst 的 slide 来清页**——
    会触发 slide1.xml 重复写、生成文件在 PowerPoint 打开即报"修复"弹窗损坏。
    始终复用模板第 1 页作封面载体，仅清空占位符文字，绝不删 sp。
    """
    return prs.slides[0]

def add_content(prs):
    """添加内容页：克隆企业模板的「品牌版式」(index 6 = 4_标题幻灯片)。
    该版式自带 logo(右上) + 顶部标题带(组合1)，内容页经版式继承自动带上品牌，
    主体区继承母版渐变底。python-pptx save() 时仍会注入 <bg><bgPr><noFill/></bgPr></bg>
    覆盖渐变，故 save 后必须调用 clean_no_fill_bg() 清除。

    ⚠️ 铁律：品牌元素在「版式层」(layout-level)，不在 slide XML 形状里——
    这是企业风"克隆模板"的正确实现：内容页复用带品牌的版式，而非空白版式(8)。
    validate_deck #10 通过检查 slide.slide_layout 的品牌形状来防静默丢失。"""
    return prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])


def clean_no_fill_bg(pptx_path, skip_slides=()):
    """后处理：清除 python-pptx 在 save() 时注入的 <cSld:bg><bgPr><noFill/></bgPr></bg>。
    该元素会覆盖企业模板的母版渐变底导致内容页全白。
    在 prs.save(path) 之后立即调用。skip_slides 为跳过的幻灯片索引（如封面）。
    """
    import zipfile
    from lxml import etree
    PP = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    AA = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    z = zipfile.ZipFile(pptx_path, 'r')
    # 找到所有内容页 slide XML
    slides = {}
    for name in z.namelist():
        if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
            idx = int(name.replace('ppt/slides/slide', '').replace('.xml', ''))
            if idx not in skip_slides:
                slides[idx] = name
    if not slides:
        z.close()
        return

    # 全部读入内存，修改后写回原路径（避免沙箱拦截 os.replace/os.remove）
    modified = False
    out_items = []
    for item in z.infolist():
        data = z.read(item.filename)
        if item.filename in slides.values():
            root = etree.fromstring(data)
            cSld = root.find('{%s}cSld' % PP)
            if cSld is not None:
                bg = cSld.find('{%s}bg' % PP)
                if bg is not None:
                    noFill = bg.find('.//{%s}noFill' % AA)
                    if noFill is not None:
                        cSld.remove(bg)
                        modified = True
                        data = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone='yes')
        out_items.append((item.filename, data))
    z.close()

    if modified:
        out = zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED)
        for fname, data in out_items:
            out.writestr(fname, data)
        out.close()

# ═══════════════════════════════════════════
# 原版 cover（逐行照抄 biz_engine.py line 179-207）
# ═══════════════════════════════════════════
def _remove_static_text(slide, target):
    """删除 slide 上文字内容等于 target 的普通文本框（非占位符）。
    例如模板烧录的「这是副标题」是固定文本框，remove_placeholders 删不掉，需单独处理。"""
    spTree = slide._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for sp in spTree.findall(qn('p:sp')):
        txBody = sp.find(qn('p:txBody'))
        if txBody is None:
            continue
        txt = ''.join(t.text or '' for t in txBody.iter(qn('a:t')))
        if txt.strip() == target:
            sp.getparent().remove(sp)

def cover(slide, title, subtitle, subtitle2='', date_str='', total=1, n=1):
    """封面（参照特来电企业风截图 v3）：
    - 顶部标识 FDE崛起：28pt Bold 白 左对齐（title 文本框，y=0.4）
    - 主标题(subtitle)：白色 36pt Bold 居中 y=2.15，视觉重心
    - 强调行(subtitle2)：品牌青 20pt 居中 y=3.25，左右引号 + 横向装饰线
    - 日期：13pt 灰色居中 y=4.35
    - QR：底中，青色细边框方框包围（见 add_qr_enterprise）
    - 占位符清空文字 + 隐藏形状（hidden=1），避免主标题重复显示"""
    # 清空占位符文字 + 隐藏占位符形状（避免显示"单击此处添加标题"提示文字导致主标题重复）
    # 保留 XML 元素不删除（避免 layout/slide 结构不一致触发修复提示）
    # 同时移除 bodyPr 内的 normAutofit/spAutoFit 元素
    for ph in list(slide.placeholders):
        try:
            tf = ph.text_frame
            tf.clear()
            bodyPr = tf._txBody.find(qn('a:bodyPr'))
            if bodyPr is not None:
                for autofit_tag in ('a:normAutofit', 'a:spAutoFit'):
                    el = bodyPr.find(qn(autofit_tag))
                    if el is not None:
                        bodyPr.remove(el)
            # 隐藏占位符形状，防止显示默认提示文字与自定义标题重复
            sp = ph._element
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    cNvPr.set('hidden', '1')
        except Exception:
            pass
    _remove_static_text(slide, '这是副标题')
    # 主标题：角色 h1（28pt 白）左对齐，文字在框内垂直居底
    tf = tb_box(slide, 0.6, 0.4, 8.3, 0.5)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    apply_role(tf, "h1", title, ENT_PALETTE, PP_ALIGN.LEFT, Pt(0))
    # 主标题：白色 hero 居中（视觉重心）；字号智能算——从 36pt 起按内容区(0.6→12.73=12.13")
    # 反解能单行容纳的最大字号，长副标题自动缩至 ≤32pt，保证居中主标题不换行（SKILL §4.2.1）
    if subtitle:
        hero_max, hero_bold = TYPOGRAPHY["hero"]
        hero_pt = fit_cover_title(subtitle, hero_max, min_pt=24)   # 占屏≤72%
        hero_tf = cover_line(slide, 0.6, 2.65, 12.13, subtitle, hero_pt,
                             ENT_PALETTE["hero"], hero_bold, PP_ALIGN.CENTER, cy=2.65)
    # 强调行：纯文字（无引号、无装饰线）；封面特有 20pt 青，沿用已定稿尺寸
    if subtitle2:
        y2 = shape_bottom_in(hero_tf) + 0.18
        sub2_tf = cover_line(slide, 0.8, y2, 11.73, subtitle2, 20, BIZ_ACCENT,
                             False, PP_ALIGN.CENTER)
    # 日期
    if date_str:
        anchor_tf = sub2_tf if subtitle2 else hero_tf
        y3 = shape_bottom_in(anchor_tf) + 0.22
        cover_line(slide, 0.8, y3, 11.73, date_str, 13, ENT_MUTE,
                   False, PP_ALIGN.CENTER)
    add_qr_enterprise(slide)   # 封面二维码：底中(6.17,5.5)+标签（docstring 既定，此前漏调用）
    add_page_number(slide, n, total, ENT_MUTE)

# ═══════════════════════════════════════════
# section_header（定稿版：沉淀自时间轴迭代验收）
# ═══════════════════════════════════════════
def section_header(slide, h1, label, sub=None, n=None, total=None):
    """内容页标题区（定稿版）：
    - H1：28pt Bold 白，x=0.46 y=0.25 w=8.6，BOTTOM anchor
    - label：18pt Bold 金 #C9A227，y=1.10 向下拉开，与 H1 不重叠
    - sub（金句）：渲染到底部居中（18pt 金），由 bottom_gold 样式承载；
      传空/不传则不渲染（其余页面金句用 bottom_gold() 单独调用）"""
    tf = tb_box(slide, 0.46, 0.25, 8.6, 0.65)
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    apply_role(tf, "h1", h1, ENT_PALETTE, PP_ALIGN.LEFT, Pt(0), first=True)
    if label:
        tf = tb_box(slide, 0.46, 1.10, 11.73, 0.45)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = apply_role(tf, "label", label, ENT_PALETTE, PP_ALIGN.LEFT, Pt(0), first=True)
        if sub:
            # 副标题与 label 同一行（label 青粗 + sub 灰），不再落页底压金句
            r2 = p.add_run(); r2.text = "    " + sub
            r2.font.size = Pt(14); r2.font.bold = False; r2.font.color.rgb = LIGHT_GRAY
            r2.font.name = '微软雅黑'; set_chinese_font(r2, '微软雅黑')
        # 用 label+sub 合并宽度定框，避免 sub 被塞进只按 label 宽的窄框而换行（企业风 P3 副标题 bug）
        fit_single_line(tf, (label + "    " + sub) if sub else label,
                        TYPOGRAPHY["label"][0], buffer=0.2, max_w=12.27)
    if n and total:
        add_page_number(slide, n, total, ENT_MUTE)

def _render_gold(slide, text):
    """底部金句绘制（不含页码）：gold 角色（16pt Bold 橙 #F79646，四风格统一 16pt；无分隔线），y=6.65 居中。"""
    tf = tb_box(slide, TL_SIDE, 6.65, TL_CARD_RIGHT - TL_SIDE, 0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    apply_role(tf, "gold", text, ENT_PALETTE, PP_ALIGN.CENTER, Pt(0), first=True)
    w = fit_single_line(tf, text, TYPOGRAPHY["gold"][0], buffer=0.2, max_w=12.13)  # 用满内容区(0.6→12.73)
    tf._parent.left = Inches((PAGE_W - w) / 2.0)   # 居中：左右均有余量、单行

# ═══════════════════════════════════════════
# two_column（原版 two_column line 242-263 直译）
# ═══════════════════════════════════════════
def two_column(slide, title=None, left=None, right=None, n=None, total=None,
               label=None, sub=None, left_note=None, right_note=None, right_flow=False):
    """双栏（企业风，委托通用 layout_two_column）：
    - 标题+副标题统一委托 section_header 渲染，保证与全站 H1(28白)/label(18青@y=1.10)
      字号与间距一致（参照 Skill 首页标题区规范），消除双栏页副标题 14pt/间距过紧的漂移
    - 卡片与 card_grid 同款（content_card）；栏标题 16pt 白粗 + 青色短下划线
    - 卡宽卡高由内容测算取本页最大值统一（通用铁律）；两卡对称居中
    - 底部青色结论 note；right_flow=True 时右栏渲染为时序流（青色阶段词高亮）"""
    if title:
        # 标题区交给 section_header：H1(28白) + label(18青@y=1.10)，与全站一致；sub 同行内联（同 section_header 约定）
        section_header(slide, title, label, sub=sub, n=n, total=total)
    else:
        if n and total:
            add_page_number(slide, n, total, ENT_MUTE)
    layout_two_column(slide, ENTERPRISE_SKIN, {"columns": [left, right]},
                      y0=1.62, y_end=6.45,
                      notes=[left_note, right_note], right_flow=right_flow, underline=True)
    # 页码已由 section_header 处理（title 存在时），避免重复落页码

# ═══════════════════════════════════════════
# card_grid — 沿用原版两栏布局风格
# 严格按 SKILL.md「2×2 卡片规则」: 标题 20pt Bold CENTER, 正文 14pt LEFT
# ═══════════════════════════════════════════
def card_grid(slide, cards, cols=2, n=None, total=None):
    """2×2 卡片：公共骨架 layout_card_grid（透明底+青边框，向带框卡约定对齐）。"""
    layout_card_grid(slide, ENTERPRISE_SKIN, cards, cols=cols, y0=1.60, y_end=6.38)
    if n and total:
        add_page_number(slide, n, total, ENT_MUTE)
def three_column(slide, title=None, label=None, cards=None, n=None, total=None, connect=False):
    """三列（企业风，委托通用 layout_three_column）：单排三卡、等宽等框高，与 2×2 一致（无下划线/无 icon/无连接线）。
    标题区复用 section_header（28pt H1 + 18pt label，与全站一致）；n/total 由 section_header 绘制页码。"""
    if title:
        section_header(slide, title, label, n=n, total=total)
    else:
        if n and total:
            add_page_number(slide, n, total, ENT_MUTE)
    layout_three_column(slide, ENTERPRISE_SKIN, cards, y0=1.65, y_end=6.45, connect=connect)
def value_loop(slide, h1, label, items, gold, n=None, total=None):
    """四大价值/闭环页：公共骨架 layout_value_grid（loop=True：四段等长循环箭头、四卡等距 0.7）。
    items 顺序固定：左上(01) → 右上(02) → 右下(03) → 左下(04)（顺时针）。"""
    section_header(slide, h1, label)
    layout_value_grid(slide, ENTERPRISE_SKIN, items, y0=1.78, y_end=6.50, loop=True)
    bottom_gold(slide, gold, n, total)
def _value_card(slide, x, y, w, h, it):
    """单张价值卡：透明底 + 亮青边框；编号(28青,左) + 标题(16白,左,同行贴编号)
       + 要点(12灰,间距0.32) + 暗色水平分隔线 + 结论(16青,居中)。"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.background()                       # 无背景填充
    card.line.color.rgb = BIZ_ACCENT            # 亮青边框
    card.line.width = Pt(1.25)
    sp = card._element
    style = sp.find(qn('p:style'))
    if style is not None: sp.remove(style)

    # 编号（28pt 青，左对齐，垂直居中于标题行）
    tf = tb_box(slide, x + 0.25, y + 0.20, 0.70, 0.55)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["num"], 28, True, BIZ_ACCENT, PP_ALIGN.LEFT, Pt(0), first=True)

    # 标题（16pt 白，左对齐，与编号同一行、垂直居中贴编号）
    tf = tb_box(slide, x + 1.0, y + 0.20, w - 1.25, 0.55)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["title"], 16, True, WHITE, PP_ALIGN.LEFT, Pt(0), first=True)

    # 要点（左对齐灰 12pt，标题→正文间距收紧至 ~0.07in，行间 0.32in）
    for j, ln in enumerate(it.get("points", [])):
        tf = tb_box(slide, x + 0.35, y + 0.82 + j * 0.32, w - 0.6, 0.32)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run_text(tf, "· " + ln, 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, Pt(0), first=True)

    # 水平分隔线（略亮蓝灰 #44586C，隐约可见、仍低于结论青高亮）
    rect(slide, x + 0.25, y + 1.52, w - 0.5, Pt(1), fill=VALUE_DIV_CLR)

    # 结论（底部居中青 16pt，记忆钩子，位于分隔线下方）
    tf = tb_box(slide, x + 0.25, y + 1.59, w - 0.5, 0.34)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["tag"], 16, True, BIZ_ACCENT, PP_ALIGN.CENTER, Pt(0), first=True)


# ───────────────────────────────────────────
# 企业风箭头约定：白 2pt（深底专属）
# 线宽 2pt = 跨风格通用（见 common.LOOP_ARROW_W）；颜色 WHITE 仅因企业风深底。
# 浅底风格(文质 / 商务)须改用深色箭头，切勿误用白。
# ───────────────────────────────────────────
def _value_arrows(slide, mx, rx, ty, by, cw, ch, L):
    """4 段等长直线箭头（企业风：白 2pt，顺时针）：01→02(右) 02→03(下) 03→04(左) 04→01(上)。
    四箭头长度均 = L（中缝），保证四卡上下左右间距全等；坐标由卡片外框精确推算。"""
    top_y = ty + ch / 2                 # 顶/底行卡中部 y
    bot_y = by + ch / 2
    l_mid = mx + cw / 2                 # 左/右列 x 中点
    r_mid = rx + cw / 2
    loops = [
        (mx + cw, top_y,  rx,      top_y),     # 顶：01→02 向右
        (r_mid,   ty + ch, r_mid,  by),        # 右：02→03 向下
        (rx,      bot_y,  mx + cw, bot_y),     # 底：03→04 向左
        (l_mid,   by,      l_mid,  ty + ch),   # 左：04→01 向上
    ]
    for (x1, y1, x2, y2) in loops:
        loop_arrow(slide, x1, y1, x2, y2, color=WHITE)   # 企业风：白 2pt（深底专属）

# ═══════════════════════════════════════════
# 底部金句（定稿版，复用 _render_gold 样式）
# ═══════════════════════════════════════════
def bottom_gold(slide, text, n=None, total=None, kw_a=()):
    """底部金句：四风格统一 16pt（企业风用品牌橙 #F79646 作 DNA，无分隔线），机制委托 common.bottom_gold。"""
    _gold(slide, text, color=ORANGE, size=16,
          page_color=ENT_MUTE, n=n, total=total, kw_a=kw_a)

# ═══════════════════════════════════════════
# definition_compare — 「定义 + 排除对比」版式（沉淀自定义对比迭代）
# 视觉语言：左宽(定义/认知) + 右窄(边界/排除) + 中部细竖线 + 底部橙金句
# 满足 SKILL.md 4.12：字号全≥12、左右非对称均衡、整体在内容区垂直居中
# ═══════════════════════════════════════════
NEG_BG   = RGBColor(0x3A, 0x46, 0x54)   # 排除项圆底（暗灰，视觉降温）
NEG_KW   = RGBColor(0x9A, 0xA5, 0xB0)   # 排除项关键词（退后灰）
NEG_D    = RGBColor(0x7C, 0x88, 0x94)   # 排除项描述（更暗）
POS_BG   = RGBColor(0x10, 0x33, 0x3E)   # 正解项底（暗青）
DIV_CLR  = RGBColor(0x2A, 0x3A, 0x4D)   # 左右分隔细线

def definition_compare(slide, h1, label, hero, hero_en, hero_cn,
                       left_bullets, left_summary,
                       compare_header, negatives, positive, gold,
                       n=None, total=None):
    """定义+排除对比：公共骨架 layout_compare（左宽+右窄圆底×badge+竖线，企业皮肤）。
    negatives: [{"kw","desc"}]；positive: 单句（骨架 kw 为空 → 框内仅一句）。"""
    section_header(slide, h1, label)
    layout_compare(slide, ENTERPRISE_SKIN, {
        "hero": hero, "hero_en": hero_en, "hero_cn": hero_cn,
        "bullets": left_bullets, "summary": left_summary,
        "right": {"header": compare_header,
                  "negatives": negatives,
                  "positive": {"kw": "", "desc": positive}},
    }, y0=1.85, y_end=6.30, x0=0.60, cw=12.13)
    _render_gold(slide, gold)
    if n and total:
        add_page_number(slide, n, total, ENT_MUTE)
def timeline(slide, items, n=None, total=None):
    """时间轴（定稿版）：
    - homePlate 卡片（暗青边框 #1A5560，无填充，切角在右）
    - 每张卡片对齐独立短虚线（不一线到底）
    - 日期 5 位 20pt 亮青蓝；正文 14pt 浅灰白；徽章 12pt 青蓝（有徽章占 2 行，左对齐）
    - 卡片高度按是否含徽章自动；卡片间距(gap)按可用垂直空间自动放大
    - 整体居中，左右留空对等（1.45in），卡片宽度 9.0in（不糊满屏）"""
    line_h = TL_BODY_SZ * TL_SPACING / 72.0
    badge_h = TL_BADGE_SZ * TL_SPACING / 72.0

    def card_h_for(it):
        return 2 * line_h + 2 * TL_PAD_Y if it.get("amt") else line_h + 2 * TL_PAD_Y

    row_heights = [card_h_for(it) for it in items]
    cnt = len(items)
    total_card_h = sum(row_heights)
    avail_h = TL_AREA_BOT - TL_AREA_TOP
    gap = (avail_h - total_card_h) / (cnt - 1) if cnt > 1 else 0
    gap = max(gap, 0.10)                       # 保底 0.10in
    y = TL_AREA_TOP + (avail_h - total_card_h - gap * (cnt - 1)) / 2

    # 轴线（细青蓝，末端箭头）
    tl = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
        Inches(TL_AXIS_X), Inches(TL_AREA_TOP - 0.1),
        Inches(TL_AXIS_X), Inches(TL_AREA_BOT + 0.05))
    tl.line.color.rgb = BIZ_TEAL; tl.line.width = Pt(2)
    ln = tl.line._get_or_add_ln()
    te = etree.SubElement(ln, qn('a:tailEnd'))
    te.set('type', 'triangle'); te.set('w', 'med'); te.set('len', 'med')

    for it, h in zip(items, row_heights):
        has_badge = bool(it.get("amt"))
        mid_y = y + h / 2

        # 节点圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            Inches(TL_AXIS_X - 0.06), Inches(mid_y - 0.06), Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = BIZ_TEAL; dot.line.fill.background()

        # 日期 5 位 20pt 加粗亮青蓝（垂直居中于整卡）
        tf = tb_box(slide, TL_DATE_X, y, TL_DATE_W, h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        apply_role(tf, "date", fmt_date_5digit(it["time"]), ENT_PALETTE,
                   PP_ALIGN.LEFT, Pt(0), first=True)

        # 卡片：homePlate 切角在右（不翻转），无填充仅暗青边框
        add_homeplate(slide, TL_CARD_X, y, TL_CARD_W, h,
                      line_color=TL_CARD_BORDER, line_w=1,
                      fill_color=None, adj_pct=TL_CARD_ADJ)

        # 每卡独立短虚线（长度 = 卡片高，x 在日期与卡片之间）
        add_dashed_line(slide, TL_DASH_X, y + 0.04, h - 0.08, TL_DASH_CLR)

        # 文字避开右侧切角
        chamfer_w = h * (TL_CARD_ADJ / 100.0)
        text_left = TL_CARD_X + chamfer_w + TL_PAD_X * 0.5
        text_w = TL_CARD_W - chamfer_w - TL_PAD_X * 1.5

        # 第 1 行：正文 14pt 浅灰白；关键词高亮（kw_teal 青 / kw_amber 金，整段去重，落实"最多2主题高亮"）
        tf = tb_box(slide, text_left, y + TL_PAD_Y, text_w, line_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        kw_t = it.get("kw_teal") or []
        kw_a = it.get("kw_amber") or []
        emph_runs(tf, it["text"], TYPOGRAPHY["body"][0], ENT_PALETTE["body"],
                  kw_t, kw_a, color_t=BIZ_TEAL, color_a=GOLD, font=FONT)

        # 第 2 行（仅徽章项）：徽章 12pt 青蓝，居左
        if has_badge:
            bx = text_left
            by = y + TL_PAD_Y + line_h
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bx), Inches(by), Inches(0.85), Inches(badge_h))
            badge.fill.solid(); badge.fill.fore_color.rgb = TL_BADGE_BG
            badge.line.color.rgb = TL_CARD_BORDER; badge.line.width = Pt(0.75)
            btf = badge.text_frame; btf.word_wrap = True
            btf.margin_left = Emu(28000); btf.margin_right = Emu(28000)
            btf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = btf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = it["amt"]
            r.font.size = Pt(TYPOGRAPHY["badge"][0]); r.font.bold = True; r.font.name = FONT
            r.font.color.rgb = ENT_PALETTE["badge"]; set_chinese_font(r)

        y += h + gap
    if n and total:
        add_page_number(slide, n, total, ENT_MUTE)


def add_qr_enterprise(slide, path=None, label="扫码关注「天戈朱」"):
    """企业风专用 QR：底中 + 标签，与署名不重叠。"""
    path = path or qr_path()
    if not os.path.exists(path):
        return
    qw = qh = 0.95
    qx, qy = 6.17, 5.5
    slide.shapes.add_picture(path, Inches(qx), Inches(qy), Inches(qw), Inches(qh))
    tf = tb_box(slide, qx - 0.3, qy + qh + 0.05, qw + 0.6, 0.55)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, label, 9, False, LIGHT_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)

import os

# ═══════════════════════════════════════════
# work_people — 「如何工作 + 谁适合」版式（沉淀自闭环+角色迭代）
# 三段式：① 反馈闭环(4 节点水平循环) ② 两类角色(Echo/Delta) ③ 三项能力
#         + 底部橙警示金句。配色/边框与 value_loop 一致(透明底 + 青边框)。
# 满足 4.12：所有卡按内容区左右 0.6 对齐、卡宽全等、字号≥12、零越界。
# ═══════════════════════════════════════════
def work_people(slide, h1, label, loop_nodes, roles, abilities,
                bg_caption=None, contrast=None, gold=None, n=None, total=None):
    """「如何工作，谁适合做？」：公共骨架 layout_loop_page（①闭环 ②双卡⇄ ③能力卡）。
    loop_nodes: [str 或 {name,sub}]；roles: [{title,desc}]；abilities: [{title,desc}]。"""
    section_header(slide, h1, label)
    content = {
        "cap1": "① 工作方式：业务现场的一体化闭环",
        "nodes": [n if isinstance(n, dict) else {"name": n} for n in loop_nodes],
        "cap2": "② Palantir 的两类角色：边听边做，边做边改",
        "cards": [{"title": it["title"], "desc": it["desc"]} for it in roles],
        "cap3": "③ 三项核心能力",
        "tail": [{"title": it["title"], "desc": it["desc"]} for it in abilities],
    }
    if contrast:
        content["contrast"] = contrast
    if bg_caption:
        content["bg_caption"] = bg_caption
    layout_loop_page(slide, ENTERPRISE_SKIN, content, y0=1.77, y_end=6.45,
                     x0=0.60, cw=12.13, card_h=0.86, tail_h=0.60)
    if gold:
        bottom_gold(slide, gold, n, total)
    elif n and total:
        add_page_number(slide, n, total, ENT_MUTE)
def _band_sub(slide, x, y, text):
    tf = tb_box(slide, x, y, 12.13, 0.30)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, text, 14, True, BIZ_ACCENT, PP_ALIGN.LEFT, Pt(0), first=True)


def _loop_node(slide, x, y, w, h, it):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.background()
    card.line.color.rgb = ENTER_CARD_LINE
    card.line.width = Pt(1.25)
    add_soft_shadow(card)
    sp = card._element
    style = sp.find(qn('p:style'))
    if style is not None: sp.remove(style)
    tf = tb_box(slide, x + 0.15, y + 0.08, w - 0.3, 0.30)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["name"], 16, True, WHITE, PP_ALIGN.CENTER, Pt(0), first=True)
    tf = tb_box(slide, x + 0.15, y + 0.40, w - 0.3, 0.26)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["sub"], 12, False, LIGHT_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)


def _loop_arrow(slide, x1, y1, x2, y2):
    """企业风闭环箭头：白 2pt（深底专属）。通用机制见 common.loop_arrow。"""
    return loop_arrow(slide, x1, y1, x2, y2, color=WHITE)


def _loop_return_u(slide, xs, nw, ny, nh):
    """U 型闭环回流：走公共机制 common.loop_return_u（起点/终点均节点底中，左右对称）。"""
    centers = [x + nw / 2 for x in xs]
    return loop_return_u(slide, centers, ny + nh, ny + nh + 0.22, color=WHITE)


def _role_card(slide, x, y, w, h, it):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.background()
    card.line.color.rgb = ENTER_CARD_LINE
    card.line.width = Pt(1.25)
    add_soft_shadow(card)
    sp = card._element
    style = sp.find(qn('p:style'))
    if style is not None: sp.remove(style)
    tf = tb_box(slide, x + 0.25, y + 0.12, w - 0.5, 0.40)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["title"], 16, True, WHITE, PP_ALIGN.LEFT, Pt(0), first=True)
    tf = tb_box(slide, x + 0.25, y + 0.50, w - 0.5, 0.32)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["desc"], 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, Pt(0), first=True)


def _ability_card(slide, x, y, w, h, it):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.background()
    card.line.color.rgb = ENTER_CARD_LINE
    card.line.width = Pt(1.25)
    add_soft_shadow(card)
    sp = card._element
    style = sp.find(qn('p:style'))
    if style is not None: sp.remove(style)
    tf = tb_box(slide, x + 0.2, y + 0.06, w - 0.4, 0.26)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["title"], 14, True, WHITE, PP_ALIGN.LEFT, Pt(0), first=True)
    tf = tb_box(slide, x + 0.2, y + 0.32, w - 0.4, 0.26)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run_text(tf, it["desc"], 12, False, LIGHT_GRAY, PP_ALIGN.LEFT, Pt(0), first=True)


# ═══════════════════════════════════════════
# 企业风「富原语」统一骨架约定（work_people / value_loop / definition_compare / transform_summary / profile_warning）
# 沉淀自迭代验收；新增富原语请对齐本骨架，避免节奏漂移：
#   1) 顶部统一 section_header（H1 28白 + label 18青=核心金句）
#   2) 多块承载时：每块带标题(青,左缘0.6)，块间距统一 0.22″，标题居顶、框高 0.72″
#   3) 内容区统一 左0.6 / 右12.73（等宽对齐），块内/块间节奏一致，不横跳极端
#   4) 闭环/箭头一律 common.loop_arrow(color=WHITE)（深底专属：白2pt）；起终点落节点底中、对称
#   5) 底部一律 bottom_gold/_render_gold（橙20pt，单行用 fit_single_line 防换行）
#   6) 对照类内容：旧态降温(NEG_BG 暗灰+✕)、新态高亮(青边框透明底+✓)，复用同一配色语言
# ═══════════════════════════════════════════
# transform_summary — 「传统研发如何转型 + 总结」版式（沉淀自转型对照迭代）
# 视觉语言：每行 = 编号(青) + 旧态(dim 灰框) → 白2pt箭头 → 新态(青边框高亮)
#           + 底部橙金句。旧态降温用 NEG_BG，新态高亮用青边框透明底（与 work_people 一致）。
# 满足 4.12：行高全等、左右对齐 0.6/12.73、字号≥12、零越界。
# ═══════════════════════════════════════════
def transform_summary(slide, h1, label, rows, gold=None, n=None, total=None):
    """「传统研发如何转型 + 总结」：公共骨架 layout_transition_rows（旧态→新态，行高自适应）。
    rows: [{"from", "to_title", "to_desc"}]"""
    section_header(slide, h1, label)
    layout_transition_rows(slide, ENTERPRISE_SKIN, {
        "cap": "三个转变",
        "rows": [{"from": it["from"], "to_title": it["to_title"], "to_lines": [it["to_desc"]]}
                 for it in rows],
    }, y0=1.77, y_end=6.30, x0=0.60, cw=12.13)
    if gold:
        bottom_gold(slide, gold, n, total)
    elif n and total:
        add_page_number(slide, n, total, ENT_MUTE)
def profile_warning(slide, h1, label, left, right, gold=None, n=None, total=None):
    """「人才画像 + 清醒提醒」：公共骨架 layout_profile_warning（非对称实心卡+两级分隔+结论竖线）。
    left:  {"title", "roles"[短词两列/长词独占行], "conclusion"}
    right: {"title", "warns"[{kw,desc}], "conclusion"}"""
    section_header(slide, h1, label)
    layout_profile_warning(slide, ENTERPRISE_SKIN, {"left": left, "right": right},
                           y0=1.85, y_end=6.30, x0=0.60, cw=12.13)
    if gold:
        _render_gold(slide, gold)
    if n and total:
        add_page_number(slide, n, total, ENT_MUTE)
